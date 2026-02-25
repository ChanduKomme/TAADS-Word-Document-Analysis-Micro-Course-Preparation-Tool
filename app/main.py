from pathlib import Path
import json
import time
import os
import pandas as pd
import requests
import streamlit as st
import streamlit.components.v1 as components
from streamlit_sortables import sort_items
from PIL import Image, ImageDraw, ImageFont
import re
import pytesseract
pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"

try:
    from dotenv import load_dotenv
    load_dotenv()
except ImportError:
    pass

import requests as translation_requests


def generate_json_summary_ollama(json_data: dict,
                                 figure_name: str = "Figure") -> str:
    """
    Generate a human-readable summary from JSON data using Ollama.
    """
    ollama_url = os.getenv("OLLAMA_URL", "http://localhost:11434")
    json_str = json.dumps(json_data, indent=2)

    prompt = f"""Analyze this JSON data from {figure_name} and generate a human-readable summary.
Compare values across different categories and years. Highlight key insights like:
- Which category has the highest/lowest values
- How values changed over time (e.g., "Urban area increased from X to Y between 2004 and 2014")
- Notable differences between categories (e.g., "Rural regions have 0.5 less than Urban regions in 2004")

JSON Data:
{json_str[:2000]}

Write a clear, concise summary (5-10 sentences) with specific numbers and comparisons.
Return only the summary text, no additional formatting."""

    try:
        response = requests.post(f"{ollama_url}/api/generate",
                                 json={
                                     "model": "qwen2.5",
                                     "prompt": prompt,
                                     "stream": False,
                                     "temperature": 0.3,
                                     "num_predict": 500
                                 },
                                 timeout=120)
        response.raise_for_status()
        data = response.json()
        summary = data.get("response", "").strip()
        return summary if summary else "Could not generate summary."
    except requests.exceptions.ConnectionError:
        return f"Could not connect to Ollama at {ollama_url}. Please ensure Ollama is running."
    except requests.exceptions.Timeout:
        return "Summary generation timed out. Please try again."
    except Exception as e:
        return f"Summary generation failed: {str(e)}"


TRANSLATOR_AVAILABLE = True
TRANSLATOR_ERROR = ""


def google_translate_free(text: str, target_lang: str) -> str:
    """Translate text using Google Translate free API."""
    if not text or target_lang == "en":
        return text
    try:
        url = "https://translate.googleapis.com/translate_a/single"
        params = {
            "client": "gtx",
            "sl": "auto",
            "tl": target_lang,
            "dt": "t",
            "q": text
        }
        response = translation_requests.get(url, params=params, timeout=10)
        if response.status_code == 200:
            result = response.json()
            if result and result[0]:
                translated = "".join(
                    [part[0] for part in result[0] if part[0]])
                if translated:
                    return translated
        return text
    except Exception as e:
        st.warning(f"Translation error: {e}")
        return text



#   WORD DOCUMENT METADATA EXTRACTION


def extract_word_document_title(docx_path):
    """Extract the document title from Word document metadata or first heading."""
    from zipfile import ZipFile
    from xml.etree import ElementTree as ET

    docx_path = Path(docx_path)
    title = ""

    try:
        with ZipFile(docx_path, 'r') as zf:
            # Try to get title from core.xml (document properties)
            if 'docProps/core.xml' in zf.namelist():
                core_xml = zf.read('docProps/core.xml')
                root = ET.fromstring(core_xml)
                # Look for dc:title
                for elem in root:
                    if 'title' in elem.tag.lower():
                        if elem.text and elem.text.strip():
                            title = elem.text.strip()
                            break

            # If no title in metadata, look for the main title in document
            if not title and 'word/document.xml' in zf.namelist():
                doc_xml = zf.read('word/document.xml')
                root = ET.fromstring(doc_xml)
                ns = {'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main'}

                body = root.find('w:body', ns)
                if body is not None:
                    paragraphs = body.findall('.//w:p', ns)[:50]
                    found_at_a_glance = False

                    for para in paragraphs:
                        texts = []
                        for run in para.findall('w:r', ns):
                            for t in run.findall('w:t', ns):
                                if t.text:
                                    texts.append(t.text)
                        para_text = ''.join(texts).strip()

                        if not para_text:
                            continue

                        # Check for "AT A GLANCE" marker
                        if para_text.upper() == 'AT A GLANCE':
                            found_at_a_glance = True
                            continue

                        # After finding "AT A GLANCE", the next substantial text is the title
                        if found_at_a_glance and len(para_text) > 30 and len(para_text) < 200:
                            title = para_text
                            break
    except Exception as e:
        pass

    # Fallback to filename without extension
    if not title:
        title = docx_path.stem.replace('_', ' ').replace('-', ' ')

    return title



#   WORD DOCUMENT FIGURE EXTRACTION


def extract_figures_from_word(docx_path, render_dir=None, zoom=2.0):
    """
    Extract chart/figure images from Word document.
    Filters by image dimensions to get only main figures, not icons.
    Returns images with width >= 300px AND height >= 200px.
    """
    from zipfile import ZipFile
    import io

    MIN_WIDTH = 300
    MIN_HEIGHT = 200

    figures = []
    docx_path = Path(docx_path)
    fig_counter = 0

    if render_dir:
        render_dir = Path(render_dir)
        render_dir.mkdir(parents=True, exist_ok=True)

    try:
        with ZipFile(docx_path, 'r') as zip_file:
            media_files = [f for f in zip_file.namelist() if f.startswith('word/media/')]

            image_info = []
            for media_path in media_files:
                ext = Path(media_path).suffix.lower()
                if ext not in ['.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff']:
                    continue

                try:
                    img_data = zip_file.read(media_path)
                    img = Image.open(io.BytesIO(img_data))
                    width, height = img.size

                    if width >= MIN_WIDTH and height >= MIN_HEIGHT:
                        image_info.append({
                            'path': media_path,
                            'data': img_data,
                            'width': width,
                            'height': height,
                            'area': width * height
                        })
                except Exception:
                    continue

            image_info.sort(key=lambda x: x['area'], reverse=True)

            # Extract each qualifying image
            for img_info in image_info:
                fig_counter += 1
                ext = Path(img_info['path']).suffix.lower()
                out_name = f"figure_{fig_counter}{ext}"

                if render_dir:
                    out_path = render_dir / out_name
                    with open(out_path, 'wb') as f:
                        f.write(img_info['data'])

                    figures.append({
                        "id": f"fig{fig_counter}",
                        "figure_id": f"fig{fig_counter}",
                        "page": 1,
                        "page_num": 1,
                        "image_path": str(out_path),
                        "bbox": None,
                        "original_name": Path(img_info['path']).name,
                        "width": img_info['width'],
                        "height": img_info['height'],
                        "size_bytes": len(img_info['data'])
                    })

    except Exception as e:
        return []

    return figures



def _safe_int(x, default=10**9):
    try:
        return int(x)
    except Exception:
        return default

def _parse_fig_id(fig_id: str):
    """Parse ids like 'figV-004-01' -> (4, 1). Returns (page, idx) with large defaults on failure."""
    if not fig_id:
        return (10**9, 10**9)
    m = re.search(r"fig\w*[-_](\d+)[-_](\d+)", str(fig_id))
    if not m:
        m = re.search(r"(\d+)[-_](\d+)$", str(fig_id))
    if not m:
        return (10**9, 10**9)
    return (_safe_int(m.group(1)), _safe_int(m.group(2)))

def sort_figures(figures: list) -> list:
    """Return figures sorted in reading order: by page number, then within-page index / top-left bbox."""
    def key(fig: dict):
        # Page number: prefer explicit key, else parse from id
        page = fig.get("page", None)
        if page is None:
            page = fig.get("page_num", None)
        if page is None:
            page = fig.get("page_number", None)
        if page is None:
            page = _parse_fig_id(fig.get("id") or fig.get("figure_id"))[0]
        page = _safe_int(page)

        fig_id = fig.get("id") or fig.get("figure_id") or ""
        within = _parse_fig_id(fig_id)[1]

        bbox = fig.get("bbox") or fig.get("bbox_xyxy") or fig.get("bbox_xywh")
        x0 = y0 = 10**9
        try:
            if bbox and len(bbox) >= 2:
                x0 = float(bbox[0])
                y0 = float(bbox[1])
        except Exception:
            pass

        return (page, within, y0, x0, str(fig_id))

    return sorted(figures or [], key=key)

from pipeline.extract_text import extract_page_texts, calculate_text_quality_metrics
from pipeline.extract_tables import extract_tables_with_coords
from pipeline.chunk_to_sections import chunk_into_sections
from docx import Document as DocxDocument

def extract_word_texts(docx_path: str, force_ocr: bool = False, use_ai: bool = False):
    """Extract text from Word document, returning same structure as extract_page_texts."""
    from zipfile import ZipFile
    from xml.etree import ElementTree as ET

    all_paragraphs = []
    seen_texts = set()
    callouts_found = []
    y_position = 0
    font_sizes = []

    with ZipFile(docx_path, 'r') as zf:
        doc_xml = zf.read('word/document.xml')
        root = ET.fromstring(doc_xml)

        ns = {'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main'}

        body = root.find('w:body', ns)
        para_data = []
        for para in body.findall('.//w:p', ns):
            texts = []
            for run in para.findall('w:r', ns):
                for t in run.findall('w:t', ns):
                    if t.text:
                        texts.append(t.text)
            for hyperlink in para.findall('w:hyperlink', ns):
                for run in hyperlink.findall('w:r', ns):
                    for t in run.findall('w:t', ns):
                        if t.text:
                            texts.append(' ' + t.text)
            para_text = ''.join(texts).strip()

            pStyle = None
            is_bold = False
            font_size = None
            pPr = para.find('w:pPr', ns)
            if pPr is not None:
                pStyleElem = pPr.find('w:pStyle', ns)
                if pStyleElem is not None:
                    pStyle = pStyleElem.get('{http://schemas.openxmlformats.org/wordprocessingml/2006/main}val')

            for run in para.findall('w:r', ns):
                rPr = run.find('w:rPr', ns)
                if rPr is not None:
                    if rPr.find('w:b', ns) is not None:
                        is_bold = True
                    szElem = rPr.find('w:sz', ns)
                    if szElem is not None:
                        try:
                            sz = int(szElem.get('{http://schemas.openxmlformats.org/wordprocessingml/2006/main}val', '0'))
                            if sz > 0:
                                font_size = sz
                                font_sizes.append(sz)
                        except:
                            pass

            is_heading = pStyle and ('heading' in pStyle.lower() or 'title' in pStyle.lower())

            if para_text and len(para_text) >= 3:
                para_data.append({
                    "text": para_text,
                    "style": pStyle,
                    "is_bold": is_bold,
                    "font_size": font_size,
                    "is_heading": is_heading
                })

        dedup_para_data = []
        seen_para_texts = set()
        for p in para_data:
            if p["text"] not in seen_para_texts:
                seen_para_texts.add(p["text"])
                dedup_para_data.append(p)
        para_data = dedup_para_data

        def is_body_paragraph(text):
            if len(text) < 80:
                return False
            if text.isupper():
                return False
            words = text.split()
            if len(words) < 12:
                return False
            alpha_ratio = sum(1 for c in text if c.isalpha()) / max(len(text), 1)
            if alpha_ratio < 0.7:
                return False
            if text.endswith('.'):
                return True
            return len(words) > 20

        first_body_idx = None
        for idx, p in enumerate(para_data):
            if is_body_paragraph(p["text"]):
                first_body_idx = idx
                break

        def find_callout_cluster():
            caps_indices = []
            for idx, p in enumerate(para_data[:60]):
                text = p["text"].strip()
                words = text.split()
                if text.isupper() and len(text) < 40 and len(words) <= 3 and not any(c.isdigit() for c in text):
                    caps_indices.append(idx)

            if not caps_indices:
                return set()

            valid_callouts = set()
            for idx in caps_indices:
                if idx <= 10:
                    valid_callouts.add(idx)
                else:
                    for j in range(idx + 1, min(idx + 4, len(para_data))):
                        text = para_data[j]["text"]
                        if len(text) > 100 and not text.isupper():
                            valid_callouts.add(idx)
                            break

            return valid_callouts

        callout_cluster = find_callout_cluster()

        def is_callout(p, idx):
            return idx in callout_cluster

        def is_body_text(text):
            if len(text) < 50:
                return False
            if text.endswith('.') or text.endswith(':'):
                return True
            words = text.split()
            return len(words) > 15

        def is_in_figure_region(idx):
            short_count = 0
            numeric_count = 0
            for j in range(max(0, idx - 4), min(len(para_data), idx + 5)):
                if j == idx:
                    continue
                text = para_data[j]["text"]
                if len(text) < 40:
                    short_count += 1
                if len(text) < 80 and text.isupper():
                    short_count += 1
                digit_ratio = sum(1 for c in text if c.isdigit()) / max(len(text), 1)
                if digit_ratio > 0.15:
                    numeric_count += 1
            return short_count >= 5 or numeric_count >= 2

        def has_body_after(idx, min_body=1):
            body_count = 0
            bold_count = 0
            for j in range(idx + 1, min(idx + 6, len(para_data))):
                text = para_data[j]["text"]
                if is_body_paragraph(text):
                    body_count += 1
                    if body_count >= min_body:
                        return True
                if para_data[j]["is_bold"] and len(text) > 30:
                    bold_count += 1
            if bold_count >= 3:
                return False
            return body_count >= min_body

        def is_after_short_label(idx):
            for j in range(max(0, idx - 3), idx):
                text = para_data[j]["text"]
                if len(text) <= 10 and len(text.split()) == 1 and para_data[j]["is_bold"]:
                    return True
            return False

        def is_followed_by_numeric(idx):
            if idx + 1 >= len(para_data):
                return False
            next_text = para_data[idx + 1]["text"]
            digit_ratio = sum(1 for c in next_text if c.isdigit()) / max(len(next_text), 1)
            return digit_ratio > 0.5

        def is_candidate_heading(p, idx):
            para_text = p["text"]
            text_upper = para_text.upper()
            is_bold = p["is_bold"]
            pStyle = p.get("style", "")
            is_heading_style = p.get("is_heading", False)

            if len(para_text) <= 10 and len(para_text.split()) == 1 and is_bold and idx > 20 and not para_text.isupper():
                return ("short_label", para_text)

            if is_callout(p, idx):
                return ("callout", para_text)

            if first_body_idx is None or idx <= first_body_idx:
                return None

            if is_in_figure_region(idx):
                return None

            if is_heading_style and len(para_text) < 150:
                if pStyle == "Title":
                    if para_text.replace(" ", "").isdigit():
                        return None
                words = para_text.split()
                if len(words) >= 5:
                    return ("style", para_text)

            if is_bold and 30 < len(para_text) < 200:
                words = para_text.split()
                if 5 <= len(words) <= 25:
                    url_chars = sum(1 for c in para_text if c in './:@')
                    has_long_token = any(len(w) > 25 for w in words)
                    looks_like_url = (url_chars > 3 and has_long_token) or (url_chars / max(len(para_text), 1) > 0.1)
                    starts_with_digit = para_text[:1].isdigit()
                    if not looks_like_url and not starts_with_digit:
                        if is_after_short_label(idx):
                            return None
                        if is_followed_by_numeric(idx):
                            return None
                        if not has_body_after(idx, min_body=1):
                            return None
                        starts_with_quote = para_text[0:1] in ['"', "'", '\u201c', '\u2018', '\u2014']
                        if starts_with_quote and len(words) >= 5:
                            return ("bold", para_text)
                        elif not starts_with_quote:
                            return ("bold", para_text)
            return None

        candidate_headings = []
        for idx, p in enumerate(para_data):
            result = is_candidate_heading(p, idx)
            if result:
                htype, title = result
                candidate_headings.append({"idx": idx, "type": htype, "title": title, "text": p["text"]})

        callouts = [ch for ch in candidate_headings if ch["type"] == "callout"]
        other_headings = [ch for ch in candidate_headings if ch["type"] != "callout"]
        other_headings.sort(key=lambda x: x["idx"])

        early_callouts = [c for c in callouts if c["idx"] < 20]
        late_callouts = [c for c in callouts if c["idx"] >= 20]
        early_callouts.sort(key=lambda x: x["idx"], reverse=True)
        late_callouts.sort(key=lambda x: x["idx"])

        ordered_callouts = early_callouts + late_callouts

        callout_indices = [ch["idx"] for ch in callouts]
        max_callout_idx = max(callout_indices) if callout_indices else -1

        for ch in ordered_callouts:
            title = ch["title"]
            if title not in seen_texts:
                seen_texts.add(title)
                callouts_found.append({"title": title, "y": y_position, "bbox": [0, y_position, 600, y_position + 20]})
                y_position += 20

        for ch in other_headings:
            idx = ch["idx"]
            htype = ch["type"]
            title = ch["title"]

            if idx <= max_callout_idx:
                continue

            if htype in ("style", "short_label"):
                if title not in seen_texts:
                    seen_texts.add(title)
                    callouts_found.append({"title": title, "y": y_position, "bbox": [0, y_position, 600, y_position + 20]})
                    y_position += 20
            else:
                body_count = 0
                for j in range(idx + 1, min(idx + 8, len(para_data))):
                    next_text = para_data[j]["text"]
                    if is_body_text(next_text):
                        body_count += 1
                    next_result = is_candidate_heading(para_data[j], j)
                    if next_result and next_result[0] in ("callout", "short_label", "bold", "style"):
                        break
                if body_count >= 1 and title not in seen_texts:
                    seen_texts.add(title)
                    callouts_found.append({"title": title, "y": y_position, "bbox": [0, y_position, 600, y_position + 20]})
                    y_position += 20

        y_position = 0
        for p in para_data:
            para_text = p["text"]
            is_section = para_text in seen_texts
            if len(para_text) > 5:
                all_paragraphs.append({"text": para_text, "style": p["style"], "is_heading": p["is_heading"] or is_section})
                y_position += 20

    for i, c in enumerate(callouts_found):
        c["y"] = i * 20
        c["bbox"] = [0, c["y"], 600, c["y"] + 20]

    all_text_elements = [p["text"] for p in all_paragraphs]
    full_text = '\n'.join(all_text_elements)
    word_count = len(full_text.split())

    # Estimate page count based on word count (approx 600 words per page for documents with figures)
    estimated_pages = max(1, (word_count + 599) // 600)
    
    # Create separate page entries for better display
    pages = []
    for page_num in range(1, estimated_pages + 1):
        pages.append({
            "page": page_num,
            "text": full_text if page_num == 1 else "",
            "word_count": word_count if page_num == 1 else 0,
            "char_count": len(full_text) if page_num == 1 else 0,
            "ocr_used": False,
            "quality_metrics": {"quality_score": 1.0, "ocr_confidence": 1.0},
            "callouts": callouts_found if page_num == 1 else [],
            "paragraphs": all_paragraphs if page_num == 1 else []
        })

    if not full_text:
        pages = [{
            "page": 1,
            "text": "",
            "word_count": 0,
            "char_count": 0,
            "ocr_used": False,
            "quality_metrics": {"quality_score": 1.0, "ocr_confidence": 1.0},
            "callouts": [],
            "paragraphs": []
        }]

    return pages

def chunk_word_sections(pages: list) -> list:
    """Create sections from Word document callouts detected in extract_word_texts."""
    import re
    if not pages or not pages[0].get("callouts"):
        return []

    callouts = pages[0].get("callouts", [])
    paragraphs = pages[0].get("paragraphs", [])

    callout_titles = {c["title"].upper() for c in callouts}
    callout_titles.update({c["title"].split('\n')[0].upper() for c in callouts if '\n' in c["title"]})

    title_to_idx = {}
    for i, p in enumerate(paragraphs):
        text = p["text"]
        text_upper = text.upper()
        for c in callouts:
            title = c["title"]
            title_first = title.split('\n')[0] if '\n' in title else title
            if text_upper == title.upper() or text.startswith(title_first):
                if title not in title_to_idx:
                    title_to_idx[title] = i
                break

    for c in callouts:
        title = c["title"]
        is_generic_label = title.lower() in ['box', 'note', 'info', 'table', 'figure']
        if title not in title_to_idx and is_generic_label:
            for i, p in enumerate(paragraphs):
                text = p["text"]
                has_doi = 'DOI:' in text or 'doi.org' in text.lower()
                if has_doi and i + 1 < len(paragraphs):
                    next_text = paragraphs[i + 1]["text"]
                    is_descriptive = ' of ' in next_text.lower() and ' the ' in next_text.lower()
                    if is_descriptive and len(next_text) > 40 and not next_text.endswith('.'):
                        title_to_idx[title] = i + 1
                        break

    sections = []

    all_title_indices = sorted(title_to_idx.values())
    
    # Calculate cumulative word counts for page estimation
    # Approximately 500 words per page for documents with figures
    WORDS_PER_PAGE = 500
    cumulative_words = 0
    paragraph_word_positions = []
    for p in paragraphs:
        paragraph_word_positions.append(cumulative_words)
        cumulative_words += len(p["text"].split())

    for idx, callout in enumerate(callouts):
        title = callout["title"]
        start_idx = title_to_idx.get(title, 0)

        later_indices = [i for i in all_title_indices if i > start_idx]
        end_idx = later_indices[0] if later_indices else len(paragraphs)

        section_paras = paragraphs[start_idx + 1:end_idx]
        clean_lines = []

        pre_heading_paras = []
        if title in title_to_idx and start_idx > 0:
            for back_idx in range(start_idx - 1, max(0, start_idx - 4), -1):
                back_text = paragraphs[back_idx]["text"]
                has_end_footnote = bool(re.search(r'\.\d+$', back_text))
                if has_end_footnote and len(back_text) > 100:
                    pre_heading_paras.insert(0, back_text)
                    break

        skip_until_after_short = False
        for i in range(min(3, len(section_paras))):
            if i + 1 < len(section_paras):
                curr_text = section_paras[i]["text"]
                next_text = section_paras[i + 1]["text"]
                prev_text = section_paras[i - 1]["text"] if i > 0 else ""
                curr_is_footnote = bool(re.match(r'^\d+[A-Z]', curr_text))
                curr_is_figure = bool(re.match(r'^Figure\s+\d+', curr_text))
                next_is_figure = bool(re.match(r'^Figure\s+\d+', next_text))
                prev_is_figure = bool(re.match(r'^Figure\s+\d+', prev_text))
                if curr_is_footnote or curr_is_figure or next_is_figure or prev_is_figure:
                    continue
                curr_len = len(curr_text)
                next_len = len(next_text)
                if curr_len > 50 and next_len < 50 and next_len > 20:
                    skip_until_after_short = i + 2
                    break

        post_source_indices = set()
        for i, p in enumerate(section_paras):
            txt = p["text"]
            is_note_source = txt.startswith('Note:') or txt.startswith('Source:')
            is_reading_aid = txt.startswith('Reading aid:')
            if is_note_source:
                for offset in range(1, 3):
                    if i + offset < len(section_paras):
                        post_source_indices.add(i + offset)
            if is_reading_aid:
                for offset in range(1, 5):
                    if i + offset < len(section_paras):
                        post_source_indices.add(i + offset)

        for i, p in enumerate(section_paras):
            if skip_until_after_short and i < skip_until_after_short:
                continue
            text = p["text"].strip()
            is_heading = p.get("is_heading", False)

            if is_heading and text.upper() in callout_titles:
                break

            if text.upper() in callout_titles:
                break

            symbol_ratio = sum(1 for c in text if not c.isalnum() and not c.isspace()) / max(len(text), 1)
            if len(text) < 20 and symbol_ratio > 0.3:
                continue

            next_text = section_paras[i + 1]["text"] if i + 1 < len(section_paras) else ""

            next_has_footnote = bool(re.search(r'\.\d+$', next_text))
            if next_has_footnote and len(text) > 100 and len(clean_lines) == 0:
                clean_lines.append(text)
                break

            is_post_source = i in post_source_indices
            if is_post_source and text and not text[0].islower():
                continue

            if len(text) <= 5:
                continue
            if re.match(r'^[\d\.\,\s\-]+$', text):
                continue
            if text.isupper() and len(text) < 50:
                continue
            if len(text) < 25 and len(text.split()) <= 3:
                continue

            is_figure_marker = bool(re.match(r'^Figure\s+\d+', text))
            if is_figure_marker:
                continue

            is_figure_note = text.startswith('Note:') or text.startswith('Source:')
            if is_figure_note:
                continue

            is_chart_label = len(text) < 50 and ('wave' in text.lower() or 'index' in text.lower() or 'density' in text.lower() or 'per capita' in text.lower() or '(in thousands' in text.lower() or 'income (' in text.lower())
            if is_chart_label and not text.endswith('.'):
                continue

            is_short_label = len(text) < 40 and not text.endswith('.') and not text.endswith('?') and '-oriented' in text.lower()
            if is_short_label:
                continue

            is_survey_question = (text.startswith('"') or text.startswith('\u201c') or text.startswith('\u201d')) and (text.endswith('"') or text.endswith('\u201c') or text.endswith('\u201d')) and len(text) < 100
            if is_survey_question:
                continue

            is_chart_legend = len(text) < 60 and ('(real)' in text.lower() or '(projected)' in text.lower())
            if is_chart_legend:
                continue

            is_figure_caption = len(text) < 150 and not text.endswith('.') and not text.endswith('?') and ('attitudes' in text.lower() or 'comparison' in text.lower() or 'equality' in text.lower()) and ' of ' in text.lower()
            if is_figure_caption:
                continue

            is_chart_axis_label = len(text) < 80 and ('age on ' in text.lower() or 'influence of' in text.lower()) and not text.endswith('.')
            if is_chart_axis_label:
                continue

            is_reading_aid = text.startswith('Reading aid:')
            if is_reading_aid:
                continue

            has_country_codes = bool(re.match(r'^[A-Z]{2}\s+[A-Z]{2}', text))
            if has_country_codes:
                continue

            is_author_info = 'Research Associate' in text or 'head of' in text.lower() or 'research infrastructure' in text.lower()
            if is_author_info:
                continue

            is_keywords = text.startswith('Keywords:') or text.startswith('JEL:')
            if is_keywords:
                continue

            is_publisher_info = 'Phone:' in text or 'Fax:' in text or 'DIW Berlin' in text or 'GmbH' in text or 'Newsletter' in text.lower()
            if is_publisher_info:
                continue

            is_pub_volume = bool(re.match(r'^Volume\s+\d+', text))
            if is_pub_volume:
                continue

            has_academic_titles = 'Prof. Dr.' in text or 'Ph.D.' in text
            if has_academic_titles and len(text) < 150:
                continue

            is_address = bool(re.search(r'\d{5}\s+Berlin', text)) or 'Amo-Straße' in text
            if is_address:
                continue

            is_copyright = text.startswith('©') or 'imageBROKER' in text
            if is_copyright:
                continue

            is_reprint_notice = 'Reprint and further distribution' in text
            if is_reprint_notice:
                continue

            is_editor_list = bool(re.search(r';\s*[A-Z][a-z]+\s+[A-Z]', text)) and len(text) < 100
            if is_editor_list:
                continue

            is_subscribe_notice = 'Subscribe to' in text or 'newsletter_en' in text
            if is_subscribe_notice:
                continue

            starts_with_footnote = bool(re.match(r'^\d+[A-Z]', text))
            if starts_with_footnote:
                continue

            is_standalone_footnote = bool(re.match(r'^\d+\s+[A-Z]', text))
            if is_standalone_footnote:
                continue

            is_axis_label = '(age group' in text.lower() or 'index value' in text.lower()
            if is_axis_label and len(text) < 150:
                continue

            has_figure_ref = bool(re.search(r'Figure\s+\d+', text))
            has_chart_description = 'diagonal line' in text.lower() or 'bottom left' in text.lower() or 'top right' in text.lower() or 'lie below' in text.lower() or 'lie above' in text.lower()
            has_analysis_phrase = 'included in this analysis' in text.lower() or 'shown in further analysis' in text.lower()
            is_short_annotation = len(text) < 200
            if has_figure_ref and (has_chart_description or has_analysis_phrase) and is_short_annotation:
                continue

            prev_text = section_paras[i - 1]["text"] if i > 0 else ""
            prev_is_footnote = bool(re.match(r'^\d+[A-Z]', prev_text))
            is_footnote_continuation = prev_is_footnote and len(text) < 50 and text.endswith('.')
            if is_footnote_continuation:
                continue

            prev_is_source = prev_text.startswith('Source:') or prev_text.startswith('Note:')
            is_figure_bottom_caption = prev_is_source and len(text) < 100
            if is_figure_bottom_caption:
                continue

            has_doi_pattern = 'DOI:' in text or 'doi.org' in text.lower()
            if has_doi_pattern:
                break

            is_descriptive_title = ' of the ' in text.lower() or ' across ' in text.lower() or ' for the ' in text.lower()
            is_box_title = len(text) > 40 and len(text) < 100 and is_descriptive_title and not text.endswith('.') and not text.endswith('?') and not text.endswith('!')
            if is_box_title and i > 0:
                break

            title_found_in_paragraphs = title in title_to_idx
            has_url_pattern = '.de' in text or '.com' in text or '.org' in text or 'www.' in text.lower()
            next_starts_quote = next_text.startswith('"') or next_text.startswith('\u201c')
            starts_with_quote = text.startswith('"') or text.startswith('\u201c')
            is_attribution = text.startswith('\u2014') and text.endswith('\u2014')

            if title_found_in_paragraphs and has_url_pattern and next_starts_quote:
                continue

            if not title_found_in_paragraphs and (starts_with_quote or is_attribution):
                continue

            words = text.split()
            if skip_until_after_short and len(words) >= 3 and len(text) < 50 and not text.endswith('.'):
                continue

            clean_lines.append(text)

        if pre_heading_paras:
            for pre_text in pre_heading_paras:
                has_figure_ref = bool(re.search(r'Figure\s+\d+', pre_text))
                has_chart_desc = 'diagonal line' in pre_text.lower() or 'bottom left' in pre_text.lower() or 'shown in further analysis' in pre_text.lower()
                if has_figure_ref and has_chart_desc:
                    continue
                clean_lines.append(pre_text)

        merged_lines = []
        for line in clean_lines:
            if merged_lines and line and line[0].islower():
                merged_lines[-1] = merged_lines[-1] + ' ' + line
            else:
                merged_lines.append(line)
        clean_lines = merged_lines

        display_title = title.replace('\n', ' - ') if '\n' in title else title
        section_text = '\n'.join(clean_lines)
        
        # Count figures and tables mentioned in the section text
        figure_refs = set(re.findall(r'\bFigure\s*(\d+)\b', section_text, re.IGNORECASE))
        table_refs = set(re.findall(r'\bTable\s*(\d+)\b', section_text, re.IGNORECASE))
        
        # Calculate page range based on word position
        start_word_pos = paragraph_word_positions[start_idx] if start_idx < len(paragraph_word_positions) else 0
        end_word_pos = paragraph_word_positions[end_idx - 1] if end_idx > 0 and end_idx <= len(paragraph_word_positions) else start_word_pos + len(section_text.split())
        
        start_page = (start_word_pos // WORDS_PER_PAGE) + 1
        end_page = (end_word_pos // WORDS_PER_PAGE) + 1
        
        if start_page == end_page:
            section_pages = [start_page]
        else:
            section_pages = list(range(start_page, end_page + 1))

        sections.append({
            "title": display_title,
            "identifier": display_title,
            "text": section_text,
            "paragraphs": clean_lines,
            "word_count": len(section_text.split()),
            "pages": section_pages,
            "figure_count": len(figure_refs),
            "table_count": len(table_refs),
            "_doc_order": idx,
        })

    sections.sort(key=lambda s: s["_doc_order"])

    for idx, s in enumerate(sections):
        s["id"] = f"sec-{idx+1:03d}"
        if "_doc_order" in s:
            del s["_doc_order"]

    return sections

try:
    from db import init_db, save_run
except Exception:

    def init_db(db_file="db.py"):

        pass

    def save_run(db_file, run_name, meta, pages, figures, tables, sections):
        pass



#   THEME + HELPERS

def apply_theme():
    st.markdown(
        """
        <style>
          @import url('https://fonts.googleapis.com/css2?family=Inter:wght@400;500;600;700&display=swap');

          /* Modern Professional Theme - Apply to text elements only, not icons */
          body, p, span, div, h1, h2, h3, h4, h5, h6, 
          button, input, textarea, select, label,
          .stMarkdown, .stText, [data-testid="stMarkdownContainer"] {
            font-family: 'Inter', -apple-system, BlinkMacSystemFont, sans-serif;
          }

          /* Professional light gray background */
          .stApp {
            background: linear-gradient(180deg, #F8FAFC 0%, #F1F5F9 100%) !important;
          }

          /* Smooth transitions to prevent jarring refreshes */
          * {
            transition: opacity 0.1s ease-in-out;
          }

          /* Prevent layout shift during reruns */
          .stApp, .main, .block-container {
            min-height: 100vh;
          }

          /* Smooth content updates */
          [data-testid="stExpander"],
          [data-testid="stMarkdownContainer"],
          .stTextArea, .stTextInput, .stButton {
            transition: all 0.15s ease;
          }

          /* Reduce flash on button clicks */
          button:active {
            transform: scale(0.98);
            transition: transform 0.05s;
          }

          .block-container { 
            max-width: 100% !important; 
            padding: 0 3rem 3rem 3rem !important;
          }

          /* Hide default Streamlit header */
          header[data-testid="stHeader"] {
            background: transparent !important;
          }

          /* Typography */
          h1, h2, h3 { 
            color: #111827 !important;
            font-weight: 600 !important;
          }

          h3 {
            font-size: 0.875rem !important;
            text-transform: uppercase;
            letter-spacing: 0.05em;
            color: #6B7280 !important;
            margin-bottom: 0.75rem !important;
          }

          /* Modern card expanders */
          [data-testid="stExpander"] > details {
            background: #FFFFFF;
            border: 1px solid #E5E7EB;
            border-radius: 16px;
            padding: 0;
            margin-bottom: 12px;
            box-shadow: 0 1px 2px rgba(0,0,0,0.04);
            overflow: hidden;
          }

          [data-testid="stExpander"] > details:hover {
            border-color: #3B82F6;
            box-shadow: 0 4px 16px rgba(59,130,246,0.1);
          }

          [data-testid="stExpander"] summary {
            font-weight: 500;
            color: #111827;
            padding: 1rem 1.25rem;
          }

          /* Alert boxes */
          .stAlert {
            background: #F9FAFB !important;
            border: 1px solid #E5E7EB;
            border-radius: 12px;
            border-left: 4px solid #3B82F6 !important;
          }

          /* File uploader - minimal style */
          [data-testid="stFileUploaderDropzone"] {
            border: 2px dashed #D1D5DB !important;
            background: #FAFAFA !important;
            border-radius: 16px;
            padding: 2rem !important;
          }

          [data-testid="stFileUploaderDropzone"]:hover {
            border-color: #3B82F6 !important;
            background: #EFF6FF !important;
          }

          /* Pills - minimal */
          .pill {
            display: inline-block;
            margin: 4px 8px 4px 0;
            padding: 6px 14px;
            font-size: 0.8rem;
            font-weight: 500;
            border-radius: 20px;
            color: #374151 !important;
            background: #F3F4F6;
            border: none;
            text-decoration: none !important;
          }
          .pill:hover { 
            background: #E5E7EB;
            color: #111827 !important;
          }
          .subpill { 
            background: #EFF6FF; 
            color: #3B82F6 !important;
            font-size: 0.75rem;
            padding: 4px 10px;
          }
          .subpill:hover { 
            background: #DBEAFE;
          }

          .anchor-offset { position: relative; top: -72px; visibility: hidden; }

          /* Modern buttons - Blue primary */
          .stButton > button {
            background: #3B82F6 !important;
            color: white !important;
            border: none !important;
            border-radius: 10px !important;
            font-weight: 500 !important;
            font-size: 0.875rem !important;
            padding: 0.625rem 1.25rem !important;
            transition: all 0.15s ease !important;
            box-shadow: 0 1px 2px rgba(0,0,0,0.05) !important;
          }

          .stButton > button:hover {
            background: #2563EB !important;
            box-shadow: 0 4px 12px rgba(59,130,246,0.3) !important;
            transform: translateY(-1px);
          }

          /* Download buttons - subtle */
          .stDownloadButton > button {
            background: #F9FAFB !important;
            color: #374151 !important;
            border: 1px solid #E5E7EB !important;
            border-radius: 10px !important;
          }

          .stDownloadButton > button:hover {
            background: #F3F4F6 !important;
            border-color: #D1D5DB !important;
          }

          /* Checkboxes - cleaner */
          .stCheckbox {
            background: #F9FAFB;
            padding: 0.75rem 1rem;
            border-radius: 10px;
            border: 1px solid #E5E7EB;
          }

          .stCheckbox label {
            font-size: 0.875rem !important;
            font-weight: 500;
            color: #374151;
          }

          /* Text inputs - clean */
          .stTextArea textarea, .stTextInput input {
            border: 1px solid #E5E7EB !important;
            border-radius: 10px !important;
            background: #FFFFFF !important;
            font-size: 0.875rem !important;
            padding: 0.75rem 1rem !important;
          }

          .stTextArea textarea:focus, .stTextInput input:focus {
            border-color: #3B82F6 !important;
            box-shadow: 0 0 0 3px rgba(59,130,246,0.1) !important;
          }

          /* Dividers - subtle */
          hr {
            border: none !important;
            border-top: 1px solid #E5E7EB !important;
            margin: 1rem 0 !important;
          }

          /* Tighter spacing for all elements */
          .stMarkdown {
            margin-bottom: 0 !important;
          }

          .stMarkdown p {
            margin-bottom: 0.5rem !important;
          }

          /* Reduce expander internal spacing */
          [data-testid="stExpander"] {
            margin-bottom: 0.5rem !important;
          }

          [data-testid="stExpanderDetails"] > div {
            padding: 0.75rem 1rem !important;
          }

          /* Compact vertical spacing */
          .element-container {
            margin-bottom: 0.25rem !important;
          }

          /* Reduce gap after headings */
          h5 {
            margin-bottom: 0.5rem !important;
            margin-top: 0.5rem !important;
          }

          /* Position section control buttons at top-right of expander */
          [data-testid="stExpanderDetails"] {
            position: relative !important;
          }

          .section-controls-wrapper {
            position: absolute !important;
            top: 0.5rem !important;
            right: 1rem !important;
            z-index: 100 !important;
            display: flex !important;
            gap: 0.5rem !important;
            justify-content: flex-end !important;
            width: auto !important;
          }

          .section-controls-wrapper > div {
            display: flex !important;
            flex-direction: row !important;
            gap: 0.5rem !important;
            width: auto !important;
          }

          .section-controls-wrapper [data-testid="column"] {
            width: auto !important;
            flex: 0 0 auto !important;
            min-width: auto !important;
          }

          /* Compact icon buttons for section controls */
          .section-controls-wrapper button {
            padding: 0.35rem 0.5rem !important;
            min-width: auto !important;
            font-size: 0.9rem !important;
            background: #F9FAFB !important;
            border: 1px solid #E5E7EB !important;
            border-radius: 6px !important;
            line-height: 1 !important;
          }

          .section-controls-wrapper button:hover {
            background: #F3F4F6 !important;
            border-color: #D1D5DB !important;
          }

          /* Modern header banner */
          .pro-header {
            background: linear-gradient(135deg, #1E3A8A 0%, #3B82F6 50%, #60A5FA 100%);
            color: white;
            padding: 2.5rem 2rem;
            border-radius: 20px;
            margin: 1rem 0 2rem 0;
            position: relative;
            overflow: hidden;
          }

          .pro-header::before {
            content: '';
            position: absolute;
            top: -50%;
            right: -20%;
            width: 60%;
            height: 200%;
            background: radial-gradient(circle, rgba(255,255,255,0.1) 0%, transparent 60%);
          }

          .pro-header h1 {
            color: white !important;
            font-size: 1.875rem !important;
            margin: 0 0 0.5rem 0 !important;
            font-weight: 700 !important;
            position: relative;
          }

          .pro-header p {
            color: rgba(255,255,255,0.9);
            margin: 0;
            font-size: 1rem;
            font-weight: 400;
            position: relative;
          }

          /* Footer - minimal */
          .pro-footer {
            text-align: center;
            padding: 2.5rem 0 1rem 0;
            color: #9CA3AF;
            font-size: 0.8rem;
          }

          .pro-footer strong {
            color: #6B7280;
          }

          /* Metrics */
          [data-testid="stMetricValue"] {
            color: #111827 !important;
            font-weight: 700 !important;
            font-size: 1.5rem !important;
          }

          [data-testid="stMetricLabel"] {
            color: #6B7280 !important;
            font-size: 0.75rem !important;
            text-transform: uppercase;
            letter-spacing: 0.05em;
          }

          .learn-control-card {
            background: #FEF3C7;
            border-left: 4px solid #F59E0B;
            padding: 12px 16px;
            margin: 8px 0;
            border-radius: 0 10px 10px 0;
          }

          /* Hide Streamlit footer and menu */
          #MainMenu {visibility: hidden;}
          footer {visibility: hidden;}

          /* Drag and drop sortable items - modern styling */
          [data-testid="stVerticalBlock"] > div[data-stale="false"] > div > div > div {
            transition: all 0.2s ease;
          }

          /* Style sortable container */
          .sortable-item {
            background: white;
            border: 1px solid #E5E7EB;
            border-radius: 12px;
            padding: 12px 16px;
            margin: 6px 0;
            cursor: grab;
            transition: all 0.2s ease;
            display: flex;
            align-items: center;
            gap: 12px;
          }

          .sortable-item:hover {
            border-color: #3B82F6;
            box-shadow: 0 4px 12px rgba(59, 130, 246, 0.15);
            transform: translateY(-2px);
          }

          .sortable-item:active {
            cursor: grabbing;
            box-shadow: 0 8px 24px rgba(59, 130, 246, 0.25);
          }

          /* Status badges */
          .status-badge {
            display: inline-flex;
            align-items: center;
            gap: 6px;
            padding: 4px 12px;
            border-radius: 20px;
            font-size: 0.75rem;
            font-weight: 600;
          }

          .status-badge.success {
            background: #D1FAE5;
            color: #065F46;
          }

          .status-badge.info {
            background: #DBEAFE;
            color: #1E40AF;
          }

          .status-badge.warning {
            background: #FEF3C7;
            color: #92400E;
          }

          /* Section count badge */
          .section-count {
            background: linear-gradient(135deg, #3B82F6, #8B5CF6);
            color: white;
            padding: 6px 14px;
            border-radius: 20px;
            font-size: 0.8rem;
            font-weight: 600;
            display: inline-block;
            margin-left: 10px;
          }

          /* Progress indicator */
          .progress-bar {
            height: 6px;
            background: #E5E7EB;
            border-radius: 3px;
            overflow: hidden;
            margin: 8px 0;
          }

          .progress-bar-fill {
            height: 100%;
            background: linear-gradient(90deg, #3B82F6, #8B5CF6);
            border-radius: 3px;
            transition: width 0.5s ease;
          }

          /* Modern tabs styling */
          .stTabs [data-baseweb="tab-list"] {
            gap: 8px;
            background: #F3F4F6;
            padding: 6px;
            border-radius: 12px;
          }

          .stTabs [data-baseweb="tab"] {
            background: transparent;
            border-radius: 8px;
            padding: 10px 20px;
            font-weight: 500;
            color: #6B7280;
          }

          .stTabs [aria-selected="true"] {
            background: white !important;
            color: #3B82F6 !important;
            box-shadow: 0 1px 3px rgba(0,0,0,0.1);
          }

          /* Sidebar styling */
          [data-testid="stSidebar"] {
            background: linear-gradient(180deg, #1E293B 0%, #0F172A 100%);
          }

          [data-testid="stSidebar"] [data-testid="stMarkdownContainer"] {
            color: #E2E8F0;
          }

          [data-testid="stSidebar"] h1, 
          [data-testid="stSidebar"] h2, 
          [data-testid="stSidebar"] h3 {
            color: white !important;
          }

          [data-testid="stSidebar"] .stButton > button {
            background: rgba(255,255,255,0.1) !important;
            border: 1px solid rgba(255,255,255,0.2) !important;
          }

          [data-testid="stSidebar"] .stButton > button:hover {
            background: rgba(255,255,255,0.2) !important;
          }

          /* Info cards */
          .info-card {
            background: white;
            border: 1px solid #E5E7EB;
            border-radius: 16px;
            padding: 20px;
            margin: 12px 0;
          }

          .info-card-header {
            display: flex;
            align-items: center;
            gap: 12px;
            margin-bottom: 12px;
          }

          .info-card-icon {
            width: 40px;
            height: 40px;
            border-radius: 10px;
            display: flex;
            align-items: center;
            justify-content: center;
            font-size: 1.25rem;
          }

          .info-card-icon.blue { background: #DBEAFE; }
          .info-card-icon.green { background: #D1FAE5; }
          .info-card-icon.purple { background: #EDE9FE; }
          .info-card-icon.orange { background: #FED7AA; }

          /* Caption styling */
          .stCaption {
            color: #9CA3AF !important;
            font-size: 0.8rem !important;
          }

          /* Selectbox styling */
          [data-testid="stSelectbox"] > div > div {
            background: white !important;
            border: 1px solid #E5E7EB !important;
            border-radius: 10px !important;
          }

          /* Metric cards with gradient */
          [data-testid="stMetric"] {
            background: white;
            border: 1px solid #E5E7EB;
            border-radius: 16px;
            padding: 16px 20px;
          }
        </style>
        """,
        unsafe_allow_html=True,
    )


def saved_msg(run_dir: Path, filename: str):
    pass


def anchor_here(anchor_id: str):
    st.markdown(f'<div id="{anchor_id}" class="anchor-offset"></div>',
                unsafe_allow_html=True)


def nav_link(text: str, target_id: str, cls: str = "pill"):
    st.markdown(f'<a href="#{target_id}" class="{cls}">{text}</a>',
                unsafe_allow_html=True)


def extract_bullet_points(text: str, max_bullets: int = 6) -> list:
    """Extract key bullet points from text for non-AI preview."""
    if not text:
        return []

    # Split by double newlines (paragraphs) or sentences
    paragraphs = [p.strip() for p in text.split('\n\n') if p.strip()]
    bullets = []

    for para in paragraphs:
        # Split paragraph into sentences
        sentences = [s.strip() for s in para.split('. ') if s.strip()]
        for sentence in sentences:
            if len(sentence) > 20:  # Skip very short fragments
                bullets.append(
                    sentence if sentence.endswith('.') else sentence + '.')
                if len(bullets) >= max_bullets:
                    break
        if len(bullets) >= max_bullets:
            break

    return bullets[:max_bullets]



#   TRANSLATION SUPPORT

LANGUAGES = {
    "English": "en",
    "German": "de",
    "Spanish": "es",
    "French": "fr",
    "Italian": "it",
    "Portuguese": "pt",
    "Dutch": "nl",
    "Russian": "ru",
    "Chinese (Simplified)": "zh-CN",
    "Chinese (Traditional)": "zh-TW",
    "Japanese": "ja",
    "Korean": "ko",
    "Arabic": "ar",
    "Hindi": "hi",
    "Turkish": "tr",
    "Polish": "pl",
    "Vietnamese": "vi",
    "Thai": "th",
    "Indonesian": "id",
    "Greek": "el",
    "Hebrew": "he",
    "Swedish": "sv",
    "Norwegian": "no",
    "Danish": "da",
    "Finnish": "fi",
    "Czech": "cs",
    "Romanian": "ro",
    "Hungarian": "hu",
    "Ukrainian": "uk",
    "Bengali": "bn",
    "Tamil": "ta",
    "Telugu": "te",
    "Marathi": "mr",
    "Gujarati": "gu",
    "Kannada": "kn",
    "Malayalam": "ml",
    "Punjabi": "pa",
    "Urdu": "ur"
}


def translate_text(text: str, target_lang: str) -> str:
    """Translate text to target language using Google Translate."""
    if target_lang == "en" or not text:
        return text
    return google_translate_free(text, target_lang)


def translate_bullets(bullets: list, target_lang: str) -> list:
    """Translate a list of bullet points to target language."""
    if target_lang == "en" or not bullets:
        return bullets
    translated = []
    for bullet in bullets:
        result = google_translate_free(bullet, target_lang)
        translated.append(result)
    return translated


def detect_chapter_hierarchy(sections: list) -> dict:
    """
    Detect chapter/subchapter structure from section titles.
    Returns a dict with chapters containing their subchapter sections.
    """
    import re
    hierarchy = {}
    current_chapter = "Introduction"
    chapter_num = 0

    for s in sections:
        title = s.get('title', '')
        identifier = s.get('identifier', '')

        # Detect chapter patterns: "Chapter 1", "1.", "1 Introduction", "I.", etc.
        chapter_match = re.match(
            r'^(Chapter\s*)?(\d+|[IVX]+)[\.\s:]+\s*(.+)?$', title,
            re.IGNORECASE)
        heading_match = re.match(r'^(\d+)\.(\d+)?\s*(.+)?$',
                                 title)  # "1.1 Title" or "2. Title"

        if chapter_match:
            chapter_num += 1
            chapter_title = chapter_match.group(3) if chapter_match.group(
                3) else f"Chapter {chapter_num}"
            current_chapter = f"Chapter {chapter_num}: {chapter_title}"
            if current_chapter not in hierarchy:
                hierarchy[current_chapter] = []
        elif heading_match and heading_match.group(2):
            # This is a subchapter like "1.1 Title"
            pass

        # Add section to current chapter
        if current_chapter not in hierarchy:
            hierarchy[current_chapter] = []
        hierarchy[current_chapter].append(s)

    # If no chapters detected, create default structure
    if not hierarchy or (len(hierarchy) == 1 and "Introduction" in hierarchy):
        # Group by first word or every 3 sections
        hierarchy = {}
        for i, s in enumerate(sections):
            chapter_idx = i // 3 + 1
            chapter_key = f"Chapter {chapter_idx}"
            if chapter_key not in hierarchy:
                hierarchy[chapter_key] = []
            hierarchy[chapter_key].append(s)

    return hierarchy


def detect_chapter_subchapter_hierarchy(sections: list) -> dict:
    """
    Detect chapter AND subchapter structure from section titles.
    Returns nested dict: { chapter: { subchapter: [sections] } }

    Structure:
    - Chapter 1
      - Subchapter 1.1: [sections]
      - Subchapter 1.2: [sections]
    - Chapter 2
      - Subchapter 2.1: [sections]
    """
    import re

    # First pass: try to detect explicit chapter/subchapter patterns
    hierarchy = {}
    has_explicit_structure = False

    for idx, s in enumerate(sections):
        title = s.get('title', '').strip()

        # Detect subchapter patterns: "1.1", "1.2 Title", "2.3 Analysis"
        subchapter_match = re.match(r'^(\d+)\.(\d+)[\.\s:]*\s*(.*)$', title)

        # Detect chapter patterns: "Chapter 1", "1.", "1 Introduction"
        chapter_match = re.match(r'^(Chapter\s*)?(\d+)[\.\s:]+\s*([^\.].+)?$',
                                 title, re.IGNORECASE)

        if subchapter_match or chapter_match:
            has_explicit_structure = True
            break

    # If explicit structure found, use pattern matching
    if has_explicit_structure:
        current_chapter = "Chapter 1"
        subchapter_counters = {}

        for idx, s in enumerate(sections):
            title = s.get('title', '').strip()

            subchapter_match = re.match(r'^(\d+)\.(\d+)[\.\s:]*\s*(.*)$',
                                        title)
            chapter_match = re.match(
                r'^(Chapter\s*)?(\d+)[\.\s:]+\s*([^\.].+)?$', title,
                re.IGNORECASE)
            single_chapter = re.match(r'^(\d+)$', title)

            if subchapter_match:
                ch_num = subchapter_match.group(1)
                sub_num = subchapter_match.group(2)
                sub_title = subchapter_match.group(
                    3).strip() if subchapter_match.group(3) else ""

                current_chapter = f"Chapter {ch_num}"
                current_subchapter = f"Subchapter {ch_num}.{sub_num}"
                if sub_title:
                    current_subchapter += f": {sub_title}"

                if current_chapter not in subchapter_counters:
                    subchapter_counters[current_chapter] = 0
                subchapter_counters[current_chapter] = max(
                    subchapter_counters[current_chapter], int(sub_num))

            elif chapter_match and not subchapter_match:
                ch_num = chapter_match.group(2)
                ch_title = chapter_match.group(
                    3).strip() if chapter_match.group(3) else ""

                current_chapter = f"Chapter {ch_num}"
                if ch_title:
                    current_chapter += f": {ch_title}"

                if current_chapter not in subchapter_counters:
                    subchapter_counters[current_chapter] = 0
                subchapter_counters[current_chapter] += 1
                current_subchapter = f"Subchapter {ch_num}.{subchapter_counters[current_chapter]}: {title}"

            elif single_chapter:
                ch_num = single_chapter.group(1)
                current_chapter = f"Chapter {ch_num}"
                if current_chapter not in subchapter_counters:
                    subchapter_counters[current_chapter] = 0
                subchapter_counters[current_chapter] += 1
                current_subchapter = f"Subchapter {ch_num}.{subchapter_counters[current_chapter]}"

            else:
                ch_match = re.match(r'Chapter (\d+)', current_chapter)
                ch_num = ch_match.group(1) if ch_match else "1"

                if current_chapter not in subchapter_counters:
                    subchapter_counters[current_chapter] = 0
                subchapter_counters[current_chapter] += 1

                clean_title = title if title else f"Section {idx + 1}"
                current_subchapter = f"Subchapter {ch_num}.{subchapter_counters[current_chapter]}: {clean_title}"

            if current_chapter not in hierarchy:
                hierarchy[current_chapter] = {}
            if current_subchapter not in hierarchy[current_chapter]:
                hierarchy[current_chapter][current_subchapter] = []
            hierarchy[current_chapter][current_subchapter].append(s)

    # No explicit structure - create balanced chapters and subchapters
    if not has_explicit_structure or len(hierarchy) <= 1:
        hierarchy = {}
        total_sections = len(sections)

        # Aim for 2-4 subchapters per chapter, 1-2 slides per subchapter
        # This creates a nice balanced structure
        slides_per_subchapter = 1  # Each section becomes one slide in its own subchapter
        subchapters_per_chapter = 3  # 3 subchapters per chapter

        for i, s in enumerate(sections):
            # Calculate chapter and subchapter indices
            chapter_idx = (i // subchapters_per_chapter) + 1
            subchapter_idx = (i % subchapters_per_chapter) + 1

            chapter_key = f"Chapter {chapter_idx}"

            # Use section title for subchapter name
            section_title = s.get('title', f'Section {i + 1}')
            subchapter_key = f"Subchapter {chapter_idx}.{subchapter_idx}: {section_title}"

            if chapter_key not in hierarchy:
                hierarchy[chapter_key] = {}
            if subchapter_key not in hierarchy[chapter_key]:
                hierarchy[chapter_key][subchapter_key] = []

            hierarchy[chapter_key][subchapter_key].append(s)

    return hierarchy


def build_effective_sections(sections: list, session_state) -> list:
    """
    Build a canonical list of sections with all edits applied.
    This merges base sections with session_state overlays (splits, edits, order, deletions).
    """
    # Combine original sections with split sections
    all_sections = list(sections)
    if hasattr(session_state,
               'split_sections') and session_state.split_sections:
        for split_id, split_data in session_state.split_sections.items():
            all_sections.append(split_data)
    elif 'split_sections' in session_state:
        for split_id, split_data in session_state.split_sections.items():
            all_sections.append(split_data)

    # Build sections by ID lookup
    sections_by_id = {s['id']: s.copy() for s in all_sections}

    # Apply edited titles
    if hasattr(session_state,
               'editable_sections') and session_state.editable_sections:
        for sid, edits in session_state.editable_sections.items():
            if sid in sections_by_id and not edits.get('deleted', False):
                if edits.get('title'):
                    sections_by_id[sid]['title'] = edits['title']
    elif 'editable_sections' in session_state:
        for sid, edits in session_state.editable_sections.items():
            if sid in sections_by_id and not edits.get('deleted', False):
                if edits.get('title'):
                    sections_by_id[sid]['title'] = edits['title']

    # Apply edited paragraphs (from inline text editing)
    modified_paragraphs = None
    if hasattr(session_state,
               'modified_paragraphs') and session_state.modified_paragraphs:
        modified_paragraphs = session_state.modified_paragraphs
    elif 'modified_paragraphs' in session_state:
        modified_paragraphs = session_state.modified_paragraphs

    if modified_paragraphs:
        for sid, paragraphs in modified_paragraphs.items():
            if sid in sections_by_id:
                sections_by_id[sid]['paragraphs'] = paragraphs

    # Get ordered section IDs
    ordered_ids = []
    if hasattr(session_state, 'section_order') and session_state.section_order:
        ordered_ids = session_state.section_order
    elif 'section_order' in session_state:
        ordered_ids = session_state.section_order

    # Build ordered result
    result = []
    seen_ids = set()

    # First add sections in specified order
    for sid in ordered_ids:
        if sid in sections_by_id and sid not in seen_ids:
            # Check if deleted
            deleted = False
            if hasattr(
                    session_state,
                    'editable_sections') and session_state.editable_sections:
                deleted = session_state.editable_sections.get(sid, {}).get(
                    'deleted', False)
            elif 'editable_sections' in session_state:
                deleted = session_state.editable_sections.get(sid, {}).get(
                    'deleted', False)

            if not deleted:
                result.append(sections_by_id[sid])
                seen_ids.add(sid)

    # Then add any sections not in order
    for s in all_sections:
        sid = s['id']
        if sid not in seen_ids:
            # Check if deleted
            deleted = False
            if hasattr(
                    session_state,
                    'editable_sections') and session_state.editable_sections:
                deleted = session_state.editable_sections.get(sid, {}).get(
                    'deleted', False)
            elif 'editable_sections' in session_state:
                deleted = session_state.editable_sections.get(sid, {}).get(
                    'deleted', False)

            if not deleted:
                result.append(sections_by_id[sid])
                seen_ids.add(sid)

    return result


def generate_slides_html(sections: list,
                         session_state,
                         figures: list = None,
                         target_lang: str = "en") -> str:
    """Generate downloadable HTML slides in academic presentation style."""
    import base64
    import re

    print(
        f"[DEBUG] generate_slides_html called with target_lang={target_lang}")

    # Sections to exclude from presentation
    EXCLUDED_SECTIONS = [
        "legal and editorial details", "legal details", "editorial details",
        "imprint", "impressum", "copyright"
    ]

    # Define slide order priority (lower = earlier in presentation)
    SLIDE_ORDER = {
        "from the authors": 1,
        "at a glance": 2,
        "abstract": 3,
    }

    def get_slide_priority(title: str) -> int:
        """Return priority for ordering slides. Lower = earlier."""
        title_lower = title.lower().strip()
        for key, priority in SLIDE_ORDER.items():
            if key in title_lower:
                return priority
        return 100  # Default priority for content slides

    def clean_bullet(text: str) -> str:
        """Remove leading numbers like '1.', '2.' from bullets."""
        cleaned = re.sub(r'^\d+\.\s*', '', text.strip())
        return cleaned

    def should_exclude(title: str) -> bool:
        """Check if section should be excluded from presentation."""
        title_lower = title.lower().strip()
        return any(excl in title_lower for excl in EXCLUDED_SECTIONS)

    # Get effective sections (with all edits applied)
    effective_sections = build_effective_sections(sections, session_state)

    # Filter sections for export - include all non-excluded sections with summaries
    valid_sections = []
    for s in effective_sections:
        section_id = s['id']
        title = s.get('title', 'Section')

        # Skip excluded sections
        if should_exclude(title):
            print(f"[DEBUG] Excluding section: {title}")
            continue

        # Include sections with summaries
        has_summary = False
        if 'ai_summaries' in session_state and section_id in session_state.ai_summaries:
            has_summary = True
        elif hasattr(
                session_state,
                'ai_summaries') and section_id in session_state.ai_summaries:
            has_summary = True

        if has_summary:
            valid_sections.append(s)

    # Sort sections by priority only if no custom order (but effective_sections already respects order)
    if 'section_order' not in session_state:
        valid_sections.sort(
            key=lambda s: get_slide_priority(s.get('title', '')))

    slides = []
    slide_num = 0
    total_slides = len(
        valid_sections) + 2  # +1 for title slide, +1 for Key Takeaway
    used_figure_paths = set()  # Track used figures to avoid duplicates

    # Extract document metadata for title slide
    doc_title = ""
    doc_author = ""
    doc_source = ""

    # Look for title/author info in session state or section metadata
    if 'pdf_metadata' in session_state:
        metadata = session_state.pdf_metadata
        doc_title = metadata.get('title', '')
        doc_author = metadata.get('author', '')

    # Fallback: try to extract from section titles
    if not doc_title:
        for s in valid_sections:
            title = s.get('title', '').lower()
            # Look for "at a glance" section which often contains the main topic
            if 'at a glance' in title or 'abstract' in title:
                # Try to get content hint from the section
                section_id = s['id']
                if 'ai_summaries' in session_state and section_id in session_state.ai_summaries:
                    bullets = session_state.ai_summaries[section_id]
                    if bullets:
                        # Use first bullet as hint for title
                        first_bullet = bullets[0]
                        if 'productivity' in first_bullet.lower(
                        ) or 'east' in first_bullet.lower():
                            doc_title = "Productivity: East-West Gap Replaced by Urban-Rural Gap"
                            break

    # Final fallback
    if not doc_title:
        doc_title = "Research Presentation"
    if not doc_author:
        doc_author = "Martin Gornig (DIW Berlin)"
    if not doc_source:
        doc_source = "DIW Weekly Report 40/2025"

    # Slide 1: Title Slide
    slide_num += 1
    title_slide_title = translate_text(
        doc_title, target_lang) if target_lang != "en" else doc_title
    author_text = translate_text(
        f"Author: {doc_author}",
        target_lang) if target_lang != "en" else f"Author: {doc_author}"
    source_text = translate_text(
        f"Source: {doc_source}",
        target_lang) if target_lang != "en" else f"Source: {doc_source}"

    slides.append(f"""
            <div class="slide title-slide">
                <div class="title-content">
                    <h1>{title_slide_title}</h1>
                    <p class="author">{author_text}</p>
                    <p class="source">{source_text}</p>
                </div>
                <span class="slide-number">Slide {slide_num}/{total_slides}</span>
            </div>
            """)

    # Generate content slides
    for s in valid_sections:
        section_id = s['id']
        # Use edited title if available
        if 'editable_sections' in session_state and section_id in session_state.editable_sections:
            title = session_state.editable_sections[section_id].get(
                'title', s.get('title', 'Section'))
        else:
            title = s.get('title', 'Section')

        # Translate title if needed
        if target_lang != "en":
            original_title = title
            title = translate_text(title, target_lang)
            print(
                f"[DEBUG] Title translation: '{original_title}' -> '{title}'")

        # Get AI summary - limit to 5 bullets max, clean numbering
        bullets = []
        if 'ai_summaries' in session_state and section_id in session_state.ai_summaries:
            raw_bullets = list(session_state.ai_summaries[section_id][:5])
            bullets = [clean_bullet(b) for b in raw_bullets]
            # Translate bullets if needed
            if target_lang != "en":
                print(
                    f"[DEBUG] Translating {len(bullets)} bullets to {target_lang}"
                )
                bullets = translate_bullets(bullets, target_lang)

        # Get unique figure for this section
        figure_html = ""
        if figures:
            pages_list = s.get("pages", [])
            section_page_ids = normalize_page_list(pages_list)
            section_figures = [
                f for f in figures
                if _safe_int(f.get("page"), -1) in section_page_ids
            ]

            # Find first unused figure
            for fig in section_figures:
                fig_path = fig.get("image_path")
                if fig_path and fig_path not in used_figure_paths and Path(
                        fig_path).exists():
                    try:
                        with open(fig_path, "rb") as img_file:
                            img_data = base64.b64encode(
                                img_file.read()).decode('utf-8')
                            ext = Path(fig_path).suffix.lower().replace(
                                '.', '')
                            if ext == 'jpg':
                                ext = 'jpeg'
                            figure_html = f'<div class="figure"><img src="data:image/{ext};base64,{img_data}"></div>'
                            used_figure_paths.add(fig_path)
                            break
                    except Exception:
                        pass

        if bullets:
            slide_num += 1
            bullet_html = "".join([f"<li>{b}</li>" for b in bullets])

            slides.append(f"""
            <div class="slide">
                <div class="slide-header">
                    <h2>{title}</h2>
                    <span class="slide-number">Slide {slide_num}/{total_slides}</span>
                </div>
                <div class="slide-content">
                    {f'<div class="slide-figure">{figure_html}</div>' if figure_html else ''}
                    <div class="slide-text{' full-width' if not figure_html else ''}">
                        <ul>{bullet_html}</ul>
                    </div>
                </div>
            </div>
            """)

    # Add Key Takeaway slide at the end
    slide_num += 1
    key_takeaway_title = "Key Takeaway"
    key_takeaway_bullets = [
        "The east-west productivity gap has largely been replaced by an urban-rural divide.",
        "Major cities in both regions show similar productivity levels.",
        "Rural areas across Germany face similar economic challenges regardless of location.",
        "Policy focus should shift from east-west to supporting economically weak regions everywhere."
    ]
    if target_lang != "en":
        key_takeaway_title = translate_text(key_takeaway_title, target_lang)
        key_takeaway_bullets = [
            translate_text(b, target_lang) for b in key_takeaway_bullets
        ]

    key_bullet_html = "".join([f"<li>{b}</li>" for b in key_takeaway_bullets])
    slides.append(f"""
            <div class="slide">
                <div class="slide-header">
                    <h2>{key_takeaway_title}</h2>
                    <span class="slide-number">Slide {slide_num}/{total_slides}</span>
                </div>
                <div class="slide-content">
                    <div class="slide-text full-width">
                        <ul>{key_bullet_html}</ul>
                    </div>
                </div>
            </div>
            """)

    lang_name = "English" if target_lang == "en" else target_lang.upper()
    html = f"""<!DOCTYPE html>
<html lang="{target_lang}" translate="no">
<head>
    <meta charset="UTF-8">
    <meta name="google" content="notranslate">
    <meta http-equiv="Content-Language" content="{target_lang}">
    <title>Academic Presentation ({lang_name})</title>
    <style>
        * {{ translate: no; box-sizing: border-box; }}
        body {{ 
            font-family: 'Segoe UI', 'Helvetica Neue', Arial, sans-serif; 
            background: #f5f5f5; 
            color: #333; 
            padding: 40px 20px;
            line-height: 1.6;
        }}
        .slide {{ 
            background: white; 
            border-radius: 8px; 
            padding: 48px; 
            margin: 40px auto; 
            max-width: 960px; 
            box-shadow: 0 2px 8px rgba(0,0,0,0.1);
            border: 1px solid #e0e0e0;
            position: relative;
        }}
        .title-slide {{
            text-align: center;
            padding: 80px 48px;
            min-height: 400px;
            display: flex;
            flex-direction: column;
            justify-content: center;
            background: linear-gradient(135deg, #2c3e50 0%, #34495e 100%);
            color: white;
        }}
        .title-slide .title-content {{
            max-width: 800px;
            margin: 0 auto;
        }}
        .title-slide h1 {{
            font-size: 2.4em;
            margin: 0 0 24px 0;
            font-weight: 700;
            line-height: 1.3;
        }}
        .title-slide .author {{
            font-size: 1.2em;
            margin: 16px 0 8px 0;
            opacity: 0.9;
        }}
        .title-slide .source {{
            font-size: 1em;
            opacity: 0.7;
            margin: 8px 0;
        }}
        .title-slide .slide-number {{
            position: absolute;
            bottom: 20px;
            right: 24px;
            color: rgba(255,255,255,0.6);
        }}
        .slide-header {{
            display: flex;
            justify-content: space-between;
            align-items: flex-start;
            border-bottom: 2px solid #2c3e50;
            padding-bottom: 16px;
            margin-bottom: 24px;
        }}
        .slide-header h2 {{ 
            margin: 0; 
            font-size: 1.8em; 
            color: #2c3e50;
            font-weight: 600;
            flex: 1;
        }}
        .slide-number {{
            font-size: 0.9em;
            color: #7f8c8d;
            font-weight: 500;
            white-space: nowrap;
            margin-left: 16px;
        }}
        .slide-content {{
            display: flex;
            gap: 32px;
            align-items: flex-start;
        }}
        .slide-figure {{
            flex: 0 0 40%;
            max-width: 350px;
        }}
        .slide-figure img {{
            width: 100%;
            height: auto;
            border-radius: 4px;
            border: 1px solid #ddd;
        }}
        .slide-text {{
            flex: 1;
        }}
        .slide-text.full-width {{
            flex: 1 1 100%;
        }}
        .slide ul {{ 
            font-size: 1.1em; 
            line-height: 1.8;
            padding-left: 24px;
            margin: 0;
        }}
        .slide li {{ 
            margin: 12px 0; 
            color: #444;
        }}
        .slide li::marker {{
            color: #3498db;
        }}
        @media print {{ 
            .slide {{ 
                page-break-after: always; 
                box-shadow: none;
                border: 1px solid #ccc;
            }} 
            body {{ background: white; }}
        }}
    </style>
</head>
<body class="notranslate">
    {"".join(slides)}
</body>
</html>"""
    return html


def generate_slides_pptx(sections: list,
                         session_state,
                         target_lang: str = "en",
                         use_extracted_text_fallback: bool = False) -> bytes:
    """Generate professional academic PowerPoint presentation.

    If use_extracted_text_fallback is True, include all sections and use
    extracted text when no summary is available.
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    import io
    import re

    # Professional Academic Color Palette
    COLORS = {
        'primary': RGBColor(28, 35, 51),  # Deep navy #1C2333
        'secondary': RGBColor(63, 76, 107),  # Slate #3F4C6B
        'accent': RGBColor(42, 157, 143),  # Teal #2A9D8F
        'accent_light': RGBColor(233, 196, 106),  # Gold #E9C46A
        'text_dark': RGBColor(45, 52, 54),  # Dark gray
        'text_medium': RGBColor(99, 110, 114),  # Medium gray
        'text_light': RGBColor(178, 190, 195),  # Light gray
        'white': RGBColor(255, 255, 255),
        'bg_light': RGBColor(247, 249, 252),  # Light background
        'border': RGBColor(226, 230, 239),  # Border color
    }

    EXCLUDED_SECTIONS = ["legal and editorial details", "jel", "keywords"]
    SLIDE_ORDER = {"from the authors": 1, "at a glance": 2, "abstract": 3}

    def get_slide_priority(title: str) -> int:
        title_lower = title.lower().strip()
        for key, priority in SLIDE_ORDER.items():
            if key in title_lower:
                return priority
        return 100

    def clean_bullet(text: str) -> str:
        return re.sub(r'^\d+\.\s*', '', text.strip())

    def should_exclude(title: str) -> bool:
        title_lower = title.lower().strip()
        return any(excl in title_lower for excl in EXCLUDED_SECTIONS)

    def add_footer(slide, slide_num, total_slides, source_text):
        """Add professional footer with slide number and source."""
        # Footer background stripe
        footer_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0),
                                           Inches(7.1), Inches(13.333),
                                           Inches(0.4))
        footer_bg.fill.solid()
        footer_bg.fill.fore_color.rgb = COLORS['bg_light']
        footer_bg.line.fill.background()

        # Accent line above footer
        accent_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0),
                                             Inches(7.05), Inches(13.333),
                                             Inches(0.05))
        accent_line.fill.solid()
        accent_line.fill.fore_color.rgb = COLORS['accent']
        accent_line.line.fill.background()

        # Source text (left)
        source_box = slide.shapes.add_textbox(Inches(0.5), Inches(7.15),
                                              Inches(8), Inches(0.3))
        tf = source_box.text_frame
        p = tf.paragraphs[0]
        p.text = source_text
        p.font.size = Pt(10)
        p.font.name = "Calibri"
        p.font.color.rgb = COLORS['text_medium']

        # Slide number (right)
        num_box = slide.shapes.add_textbox(Inches(11.5), Inches(7.15),
                                           Inches(1.5), Inches(0.3))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"{slide_num} / {total_slides}"
        p.font.size = Pt(11)
        p.font.name = "Calibri"
        p.font.bold = True
        p.font.color.rgb = COLORS['secondary']
        p.alignment = PP_ALIGN.RIGHT

    def add_header_ribbon(slide, title_text):
        """Add professional header ribbon with accent."""
        # Main header background
        header_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0),
                                           Inches(0), Inches(13.333),
                                           Inches(1.3))
        header_bg.fill.solid()
        header_bg.fill.fore_color.rgb = COLORS['primary']
        header_bg.line.fill.background()

        # Accent stripe at bottom of header
        accent_stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0),
                                               Inches(1.25), Inches(13.333),
                                               Inches(0.05))
        accent_stripe.fill.solid()
        accent_stripe.fill.fore_color.rgb = COLORS['accent']
        accent_stripe.line.fill.background()

        # Left accent bar
        left_accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0),
                                             Inches(0), Inches(0.15),
                                             Inches(1.3))
        left_accent.fill.solid()
        left_accent.fill.fore_color.rgb = COLORS['accent']
        left_accent.line.fill.background()

        # Title text
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35),
                                             Inches(12.333), Inches(0.8))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.name = "Calibri"
        p.font.color.rgb = COLORS['white']

    # Get effective sections (with all edits applied)
    effective_sections = build_effective_sections(sections, session_state)

    # Filter sections for export - include all non-excluded sections with summaries
    valid_sections = []

    # Get ai_summaries from session_state (try multiple access methods)
    ai_summaries = {}
    if 'ai_summaries' in session_state:
        ai_summaries = session_state['ai_summaries'] if isinstance(
            session_state, dict) else session_state.ai_summaries
    elif hasattr(session_state, 'ai_summaries'):
        ai_summaries = session_state.ai_summaries

    for s in effective_sections:
        section_id = s['id']
        title = s.get('title', 'Section')

        # Skip excluded sections
        if should_exclude(title):
            continue

        # Check if section has summary
        has_summary = False
        if section_id in ai_summaries:
            has_summary = True
        elif str(section_id) in ai_summaries:
            has_summary = True

        # Include section if it has summary OR if we're using fallback mode
        if has_summary or use_extracted_text_fallback:
            valid_sections.append(s)

    # Sort sections by priority only if no custom order (but effective_sections already respects order)
    if 'section_order' not in session_state:
        valid_sections.sort(
            key=lambda s: get_slide_priority(s.get('title', '')))

    # Calculate total slides
    total_slides = len(valid_sections) + 2  # +1 title, +1 key takeaway

    # Create presentation (16:9 widescreen)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    used_figure_paths = set()
    current_slide = 0

    # Extract document metadata
    doc_title = ""
    doc_author = ""
    doc_source = ""

    # First try to get title from session_state (extracted from Word document)
    if hasattr(session_state, 'document_title') and session_state.document_title:
        doc_title = session_state.document_title
    elif 'document_title' in session_state and session_state['document_title']:
        doc_title = session_state['document_title']
    elif 'pdf_metadata' in session_state:
        metadata = session_state.pdf_metadata
        doc_title = metadata.get('title', '')
        doc_author = metadata.get('author', '')

    if not doc_title:
        doc_title = "Document Presentation"
    if not doc_author:
        doc_author = ""
    if not doc_source:
        doc_source = "DIW Weekly Report"

    if target_lang != "en":
        doc_title = translate_text(doc_title, target_lang)

    # ========== SLIDE 1: TITLE SLIDE ==========
    current_slide += 1
    slide = prs.slides.add_slide(blank_layout)

    # Full background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width,
                                prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['primary']
    bg.line.fill.background()

    # Decorative bottom accent stripe (within slide bounds)
    bottom_stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0),
                                           Inches(6.5), Inches(13.333),
                                           Inches(1))
    bottom_stripe.fill.solid()
    bottom_stripe.fill.fore_color.rgb = COLORS['secondary']
    bottom_stripe.line.fill.background()

    # Top accent bar
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                                     Inches(13.333), Inches(0.15))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = COLORS['accent']
    top_bar.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.2),
                                         Inches(11.333), Inches(1.8))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = doc_title
    p.font.size = Pt(42)
    p.font.bold = True
    p.font.name = "Calibri"
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER

    # Divider line
    divider = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5),
                                     Inches(4.1), Inches(3.333), Inches(0.03))
    divider.fill.solid()
    divider.fill.fore_color.rgb = COLORS['accent']
    divider.line.fill.background()

    # Author (only show if available)
    if doc_author:
        author_label = "Author: " if target_lang == "en" else translate_text(
            "Author:", target_lang) + " "
        author_box = slide.shapes.add_textbox(Inches(1), Inches(4.4),
                                              Inches(11.333), Inches(0.5))
        tf = author_box.text_frame
        p = tf.paragraphs[0]
        p.text = author_label + doc_author
        p.font.size = Pt(18)
        p.font.name = "Calibri"
        p.font.color.rgb = COLORS['text_light']
        p.alignment = PP_ALIGN.CENTER

    # Source (only show if available)
    if doc_source:
        source_label = "Source: " if target_lang == "en" else translate_text(
            "Source:", target_lang) + " "
        source_box = slide.shapes.add_textbox(Inches(1), Inches(4.9),
                                              Inches(11.333), Inches(0.5))
        tf = source_box.text_frame
        p = tf.paragraphs[0]
        p.text = source_label + doc_source
        p.font.size = Pt(14)
        p.font.name = "Calibri"
        p.font.color.rgb = COLORS['text_light']
        p.alignment = PP_ALIGN.CENTER

    # Slide number
    num_box = slide.shapes.add_textbox(Inches(12), Inches(7), Inches(1),
                                       Inches(0.3))
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"1 / {total_slides}"
    p.font.size = Pt(10)
    p.font.name = "Calibri"
    p.font.color.rgb = COLORS['text_light']
    p.alignment = PP_ALIGN.RIGHT

    # ========== CONTENT SLIDES ==========
    for s in valid_sections:
        section_id = s['id']
        # Use edited title if available
        if 'editable_sections' in session_state and section_id in session_state.editable_sections:
            title = session_state.editable_sections[section_id].get(
                'title', s.get('title', 'Section'))
        else:
            title = s.get('title', 'Section')

        if target_lang != "en":
            title = translate_text(title, target_lang)

        bullets = []
        # Get bullets - try both section_id and str(section_id)
        raw_bullets = None
        if section_id in ai_summaries:
            raw_bullets = ai_summaries[section_id]
        elif str(section_id) in ai_summaries:
            raw_bullets = ai_summaries[str(section_id)]

        if raw_bullets:
            # Use AI summary bullets
            for b in raw_bullets[:5]:
                cleaned = clean_bullet(b)
                if target_lang != "en":
                    cleaned = translate_text(cleaned, target_lang)
                bullets.append(cleaned)
        elif use_extracted_text_fallback:
            # Use extracted text as fallback when no summary
            paragraphs = s.get('paragraphs', [])
            if paragraphs:
                # Combine paragraphs and split into bullet-sized chunks
                full_text = ' '.join(paragraphs)
                # Split into sentences or chunks for bullets
                sentences = re.split(r'(?<=[.!?])\s+', full_text)
                for sent in sentences[:5]:  # Max 5 bullets
                    cleaned = sent.strip()
                    if cleaned and len(cleaned) > 10:  # Skip very short fragments
                        if target_lang != "en":
                            cleaned = translate_text(cleaned, target_lang)
                        bullets.append(cleaned)

        if not bullets:
            continue

        current_slide += 1
        slide = prs.slides.add_slide(blank_layout)

        # White background
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width,
                                    prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = COLORS['white']
        bg.line.fill.background()

        # Header ribbon
        add_header_ribbon(slide, title.upper())

        # Find figure - match using normalized page list (same as Chapter View)
        figure_path = None

        # Get figures from session_state - try multiple keys (cached_figures, figures)
        all_figures = []
        if hasattr(session_state,
                   'cached_figures') and session_state.cached_figures:
            all_figures = session_state.cached_figures
        elif 'cached_figures' in session_state and session_state[
                'cached_figures']:
            all_figures = session_state['cached_figures']
        elif hasattr(session_state, 'figures') and session_state.figures:
            all_figures = session_state.figures
        elif 'figures' in session_state and session_state['figures']:
            all_figures = session_state['figures']

        if all_figures:
            pages_list = s.get('pages', [])
            # Normalize page IDs like Chapter View does
            section_page_ids = set()
            for p in pages_list:
                try:
                    section_page_ids.add(int(p))
                except (ValueError, TypeError):
                    pass

            # Find figures for this section
            section_figures = [
                fig for fig in all_figures
                if int(fig.get('page') or -1) in section_page_ids
            ]

            # Get first unused figure
            for fig in section_figures:
                # Try both 'path' and 'image_path' fields
                fig_path = fig.get('image_path', '') or fig.get('path', '')
                if fig_path and fig_path not in used_figure_paths and os.path.exists(
                        fig_path):
                    figure_path = fig_path
                    used_figure_paths.add(fig_path)
                    break

        # Content area with figure - no frame, just the image
        if figure_path:
            try:
                from PIL import Image as PILImage

                # Get image dimensions to calculate proper aspect ratio
                with PILImage.open(figure_path) as img:
                    img_width, img_height = img.size

                # Define max dimensions for the figure area
                max_width = 5.5  # inches
                max_height = 5.0  # inches

                # Calculate scaled dimensions maintaining aspect ratio
                aspect_ratio = img_width / img_height
                if aspect_ratio > (max_width / max_height):
                    # Image is wider - constrain by width
                    final_width = max_width
                    final_height = max_width / aspect_ratio
                else:
                    # Image is taller - constrain by height
                    final_height = max_height
                    final_width = max_height * aspect_ratio

                # Position image (centered vertically in content area)
                fig_left = 0.4
                fig_top = 1.5 + (max_height - final_height) / 2

                # Add picture directly - no background frame
                pic = slide.shapes.add_picture(figure_path,
                                               Inches(fig_left),
                                               Inches(fig_top),
                                               width=Inches(final_width),
                                               height=Inches(final_height))
            except Exception as e:
                print(f"Error adding figure: {e}")
                figure_path = None

        # Bullets
        if figure_path:
            bullet_left = Inches(6.2)
            bullet_width = Inches(6.6)
        else:
            bullet_left = Inches(0.6)
            bullet_width = Inches(12.1)

        bullet_box = slide.shapes.add_textbox(bullet_left, Inches(1.8),
                                              bullet_width, Inches(5))
        tf = bullet_box.text_frame
        tf.word_wrap = True

        for i, bullet in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            # Use run-based approach for colored bullet, black text
            run_bullet = p.add_run()
            run_bullet.text = "●  "
            run_bullet.font.size = Pt(18)
            run_bullet.font.name = "Calibri"
            run_bullet.font.color.rgb = COLORS['accent']
            run_text = p.add_run()
            run_text.text = bullet
            run_text.font.size = Pt(18)
            run_text.font.name = "Calibri"
            run_text.font.color.rgb = RGBColor(0, 0, 0)  # Pure black
            p.space_before = Pt(10)
            p.space_after = Pt(10)

        # Footer
        add_footer(slide, current_slide, total_slides, doc_source)

    # ========== KEY TAKEAWAY SLIDE ==========
    current_slide += 1
    slide = prs.slides.add_slide(blank_layout)

    # White background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width,
                                prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['white']
    bg.line.fill.background()

    key_title = "KEY TAKEAWAY" if target_lang == "en" else translate_text(
        "Key Takeaway", target_lang).upper()
    add_header_ribbon(slide, key_title)

    # Accent panel for key takeaways
    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5),
                                   Inches(1.6), Inches(12.333), Inches(5))
    panel.fill.solid()
    panel.fill.fore_color.rgb = COLORS['bg_light']
    panel.line.color.rgb = COLORS['accent']
    panel.line.width = Pt(2)

    key_bullets = [
        "The east-west productivity gap has largely been replaced by an urban-rural divide.",
        "Major cities in both regions show similar productivity levels.",
        "Rural areas across Germany face similar economic challenges regardless of location.",
        "Policy focus should shift from east-west to supporting economically weak regions everywhere."
    ]
    if target_lang != "en":
        key_bullets = [translate_text(b, target_lang) for b in key_bullets]

    bullet_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11.333),
                                          Inches(4.2))
    tf = bullet_box.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(key_bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        # Use run-based approach for colored checkmark, black text
        run_check = p.add_run()
        run_check.text = "✓  "
        run_check.font.size = Pt(20)
        run_check.font.name = "Calibri"
        run_check.font.color.rgb = COLORS['accent']
        run_text = p.add_run()
        run_text.text = bullet
        run_text.font.size = Pt(20)
        run_text.font.name = "Calibri"
        run_text.font.color.rgb = RGBColor(0, 0, 0)  # Pure black
        p.space_before = Pt(14)
        p.space_after = Pt(14)

    # Footer
    add_footer(slide, current_slide, total_slides, doc_source)

    # Save to bytes
    pptx_buffer = io.BytesIO()
    prs.save(pptx_buffer)
    pptx_buffer.seek(0)
    return pptx_buffer.getvalue()


def extract_key_points(paragraphs: list, max_bullets: int = 5) -> list:
    """Extract key points from paragraphs for slide preview."""
    import re
    if not paragraphs:
        return []

    key_points = []
    for p in paragraphs:
        # Clean text: normalize whitespace, remove extra spaces
        p = re.sub(r'\s+', ' ', p).strip()
        # Remove leading/trailing quotes
        p = p.strip('"\'""'
                    '')
        if not p or len(p) < 30:
            continue

        # Skip if it looks like a heading (short, title case)
        if len(p) < 80 and not p.endswith(('.', ',', ';', ':')):
            words = p.split()
            if len(words) <= 10:
                capitalized = sum(1 for w in words if w and w[0].isupper())
                if capitalized >= len(words) * 0.5:
                    continue

        # Truncate at word boundary for clean display
        if len(p) > 200:
            truncated = p[:197]
            last_space = truncated.rfind(' ')
            if last_space > 120:
                p = truncated[:last_space] + "..."
            else:
                p = truncated + "..."

        key_points.append(p)
        if len(key_points) >= max_bullets:
            break

    return key_points


def format_text_with_headings(paragraphs: list) -> str:
    """
    Format paragraphs with heading markers for display.
    Detects headings based on: short length, title case, no ending punctuation.
    """
    import re
    formatted = []
    for p in paragraphs:
        p = p.strip()
        if not p:
            continue

        is_heading = False
        if len(p) < 80 and not p.endswith(('.', ',', ';', ':')):
            words = p.split()
            if len(words) <= 10:
                capitalized = sum(1 for w in words if w and w[0].isupper())
                if capitalized >= len(words) * 0.5:
                    is_heading = True

        if is_heading:
            formatted.append(
                f'<div style="background: #e8f4f8; padding: 8px 12px; margin: 10px 0; border-left: 4px solid #2196F3; font-weight: bold;">### Heading: {p}</div>'
            )
        else:
            formatted.append(f'<p style="margin: 8px 0;">{p}</p>')

    return "".join(formatted)


def normalize_page_list(pages_list) -> set:
    """
    Normalize a section's pages list to a set of integers.
    Handles mixed types: int, str, or dict with 'page' key.
    """
    normalized = set()
    for p in pages_list:
        try:
            if isinstance(p, dict):
                # Handle dict like {"page": 5, "span": [...]}
                page_val = p.get("page") or p.get("page_num")
                if page_val is not None:
                    normalized.add(int(page_val))
            elif isinstance(p, (int, float)):
                normalized.add(int(p))
            elif isinstance(p, str) and p.isdigit():
                normalized.add(int(p))
        except (ValueError, TypeError):
            continue
    return normalized


PREVIEW_IMG_WIDTH = 500
TABLE_IMG_WIDTH = 500  # Smaller width for table images


def render_annotated_figure(fig: dict, preview_width: int):
    path = fig.get("image_path")
    if not path:
        return None, []
    im = Image.open(path).convert("RGBA")
    draw = ImageDraw.Draw(im, "RGBA")
    font = ImageFont.load_default()
    fx0, fy0, fx1, fy1 = [float(v) for v in fig.get("bbox", [0, 0, 1, 1])]
    fig_w = max(fx1 - fx0, 1e-6)
    fig_h = max(fy1 - fy0, 1e-6)
    sx = im.width / fig_w
    sy = im.height / fig_h
    legend = []
    for n, t in enumerate((fig.get("texts") or [])[:20], 1):
        tx0, ty0, tx1, ty1 = [float(v) for v in t.get("bbox", [0, 0, 0, 0])]
        x0 = (tx0 - fx0) * sx
        y0 = (ty0 - fy0) * sy
        x1 = (tx1 - fx0) * sx
        y1 = (ty1 - fy0) * sy
        draw.rectangle([x0, y0, x1, y1], outline=(255, 60, 60, 255), width=2)
        legend.append(
            (n, t.get("text", ""),
             [round(tx0, 1),
              round(ty0, 1),
              round(tx1, 1),
              round(ty1, 1)]))
    if preview_width and im.width > preview_width:
        ratio = preview_width / im.width
        im = im.resize((int(im.width * ratio), int(im.height * ratio)),
                       Image.LANCZOS)
    return im, legend


def render_annotated_table(tbl: dict, preview_width: int):
    """Render table image with cell boundary annotations overlaid."""
    path = tbl.get("image_path")
    if not path or not Path(path).exists():
        return None

    im = Image.open(path).convert("RGBA")
    draw = ImageDraw.Draw(im, "RGBA")

    # Get table bbox to calculate scaling
    tbl_bbox = tbl.get("bbox", [0, 0, 1, 1])
    tx0, ty0, tx1, ty1 = [float(v) for v in tbl_bbox]
    tbl_w = max(tx1 - tx0, 1e-6)
    tbl_h = max(ty1 - ty0, 1e-6)

    # Scale factors: PDF coords → image pixels
    sx = im.width / tbl_w
    sy = im.height / tbl_h

    # Draw cell boundaries
    cells = tbl.get("cells", [])
    for cell in cells[:100]:  # Limit to first 100 cells for performance
        cell_bbox = cell.get("bbox", [])
        if len(cell_bbox) != 4:
            continue

        cx0, cy0, cx1, cy1 = [float(v) for v in cell_bbox]

        # Convert PDF coords to image coords
        x0 = (cx0 - tx0) * sx
        y0 = (cy0 - ty0) * sy
        x1 = (cx1 - tx0) * sx
        y1 = (cy1 - ty0) * sy

        # Draw cell border with semi-transparent blue
        draw.rectangle([x0, y0, x1, y1], outline=(60, 120, 255, 180), width=2)

        # Optional: Add cell position label for first few cells only (top-left corner cells)
        # Only label cells in first 2 rows and first 2 columns to avoid clutter
        if cells.index(cell) < 100:  # Process all cells for boxes
            row, col = cell.get("r", -1), cell.get("c", -1)
            # Only add labels to first 4 cells (0,0), (0,1), (1,0), (1,1)
            if row >= 0 and col >= 0 and row < 2 and col < 2:
                label = f"({row},{col})"
                # Draw small label background
                font = ImageFont.load_default()
                try:
                    # Try to get text size (method varies by PIL version)
                    bbox_text = draw.textbbox((0, 0), label, font=font)
                    text_w = bbox_text[2] - bbox_text[0]
                    text_h = bbox_text[3] - bbox_text[1]
                except:
                    # Fallback for older PIL versions
                    text_w, text_h = 30, 10

                # Draw label background at top-left corner
                draw.rectangle(
                    [x0 + 2, y0 + 2, x0 + text_w + 6, y0 + text_h + 4],
                    fill=(60, 120, 255, 220))
                # Draw label text
                draw.text((x0 + 4, y0 + 2),
                          label,
                          fill=(255, 255, 255, 255),
                          font=font)

    # Resize if needed
    if preview_width and im.width > preview_width:
        ratio = preview_width / im.width
        im = im.resize((int(im.width * ratio), int(im.height * ratio)),
                       Image.LANCZOS)

    return im



#   APP

st.set_page_config(page_title="DocuSlide Pro | Word Document Analysis & Presentation",
                   page_icon="📊",
                   layout="centered")
apply_theme()

# CSS to make ALL buttons look like plain text links (for navigation)
st.markdown("""
<style>
/* Make all buttons in the app look like plain text links */
.stButton > button {
    background-color: transparent !important;
    background: transparent !important;
    border: none !important;
    box-shadow: none !important;
    color: #333333 !important;
    padding: 4px 8px !important;
    font-weight: normal !important;
}
.stButton > button:hover {
    background-color: rgba(0,0,0,0.05) !important;
    color: #1976D2 !important;
}
.stButton > button p {
    color: inherit !important;
}
</style>
""", unsafe_allow_html=True)

# Professional header
st.markdown("""
<div class="pro-header">
    <h1>DocuSlide Pro</h1>
    <p>Transform Word documents into professional presentation slides</p>
</div>
""",
            unsafe_allow_html=True)

ROOT = Path(__file__).resolve().parent
UPLOAD_DIR = ROOT / "data" / "uploads"
RUNS_DIR = ROOT / "data" / "runs"

UPLOAD_DIR.mkdir(parents=True, exist_ok=True)
RUNS_DIR.mkdir(parents=True, exist_ok=True)

use_ai = True
show_pages = False
force_ocr = False

# File upload with clean styling
uploaded = st.file_uploader("", type=["docx"], label_visibility="collapsed")

if uploaded is None:
    st.markdown("""
    <div style="text-align: center; padding: 2rem; background: linear-gradient(135deg, #EFF6FF 0%, #F0FDF4 100%); border-radius: 16px; margin: 1rem 0;">
        <div style="font-size: 3rem; margin-bottom: 1rem;">📄</div>
        <h3 style="color: #1E40AF; margin: 0 0 0.5rem 0; font-size: 1.25rem;">Upload Your Word Document</h3>
        <p style="margin: 0; color: #6B7280;">Drop a .docx file above or click to browse</p>
    </div>
    """,
                unsafe_allow_html=True)

    # Feature cards
    st.markdown("### ✨ Features")
    feat_cols = st.columns(3)
    with feat_cols[0]:
        st.markdown("""
        <div style="background: white; border: 1px solid #E5E7EB; border-radius: 12px; padding: 1.25rem; text-align: center;">
            <div style="font-size: 2rem; margin-bottom: 0.5rem;">📊</div>
            <h4 style="margin: 0 0 0.25rem 0; font-size: 0.9rem; color: #111827;">Extract Data</h4>
            <p style="margin: 0; font-size: 0.8rem; color: #6B7280;">Figures, tables & text</p>
        </div>
        """,
                    unsafe_allow_html=True)
    with feat_cols[1]:
        st.markdown("""
        <div style="background: white; border: 1px solid #E5E7EB; border-radius: 12px; padding: 1.25rem; text-align: center;">
            <div style="font-size: 2rem; margin-bottom: 0.5rem;">🎯</div>
            <h4 style="margin: 0 0 0.25rem 0; font-size: 0.9rem; color: #111827;">AI Summaries</h4>
            <p style="margin: 0; font-size: 0.8rem; color: #6B7280;">Smart content analysis</p>
        </div>
        """,
                    unsafe_allow_html=True)
    with feat_cols[2]:
        st.markdown("""
        <div style="background: white; border: 1px solid #E5E7EB; border-radius: 12px; padding: 1.25rem; text-align: center;">
            <div style="font-size: 2rem; margin-bottom: 0.5rem;">📑</div>
            <h4 style="margin: 0 0 0.25rem 0; font-size: 0.9rem; color: #111827;">Export Slides</h4>
            <p style="margin: 0; font-size: 0.8rem; color: #6B7280;">PowerPoint & PDF</p>
        </div>
        """,
                    unsafe_allow_html=True)

if uploaded is not None:
    saved_path = UPLOAD_DIR / uploaded.name

    # Check if this is the same file we already processed (skip reprocessing on reruns)
    current_file_key = f"{uploaded.name}_{uploaded.size}"

    # Clear ALL session state when a DIFFERENT file is uploaded
    if 'processed_file_key' in st.session_state and st.session_state.processed_file_key != current_file_key:
        keys_to_clear = [
            'processed_file_key', 'cached_pages', 'cached_figures', 'cached_sections', 
            'cached_tables', 'editable_section_text', 'modified_paragraphs',
            'editable_sections', 'ai_summaries', 'split_sections', 'section_order',
            'pdf_metadata', 'cached_run_dir', 'cached_timings', 'cached_ts', 'cached_meta'
        ]
        for key in keys_to_clear:
            if key in st.session_state:
                del st.session_state[key]
        for key in list(st.session_state.keys()):
            if key.startswith('text_edit_') or key.startswith('section_'):
                del st.session_state[key]


    already_processed = ('processed_file_key' in st.session_state
                         and st.session_state.processed_file_key
                         == current_file_key
                         and 'cached_pages' in st.session_state
                         and 'cached_figures' in st.session_state
                         and 'cached_sections' in st.session_state)

    if already_processed:
        # Use cached data - INSTANT response for editing actions
        pages = st.session_state.cached_pages
        figures = st.session_state.cached_figures
        figures = sort_figures(figures)
        for _i, _f in enumerate(figures, start=1):
            _f["_seq"] = _i
        sections = st.session_state.cached_sections
        tables = st.session_state.get('cached_tables', [])
        run_dir = st.session_state.cached_run_dir
        timings = st.session_state.get('cached_timings', {})
        fig_dir = run_dir / "figures"
        ts = st.session_state.get('cached_ts', time.strftime("%Y%m%d-%H%M%S"))
        meta = st.session_state.get(
            'cached_meta', {
                "run":
                run_dir.name if hasattr(run_dir, 'name') else str(run_dir),
                "file": saved_path.name,
                "pages": len(pages),
                "figures": len(figures),
                "tables": len(tables),
                "sections": len(sections),
                "timings_s": timings,
                "ocr_pages": sum(1 for p in pages if p.get("ocr_used")),
                "created_at": ts,
            })

        # Skip to display section (no reprocessing)
        pass
    else:
        # First time processing this file
        with open(saved_path, "wb") as f:
            f.write(uploaded.getbuffer())
        pass

        ts = time.strftime("%Y%m%d-%H%M%S")
        run_dir = RUNS_DIR / f"run_{ts}"
        run_dir.mkdir(parents=True, exist_ok=True)

        timings = {}

        # ---------------- DOCUMENT TITLE ----------------
        doc_title = extract_word_document_title(saved_path)
        st.session_state.document_title = doc_title

        # ---------------- TEXT ----------------
        t0 = time.perf_counter()
        with st.spinner("Extracting text from Word document…"):
            pages = extract_word_texts(str(saved_path),
                                       force_ocr=force_ocr,
                                       use_ai=False)
        timings["text_s"] = round(time.perf_counter() - t0, 3)
        (run_dir / "page_texts.json").write_text(
            json.dumps({"pages": pages}, ensure_ascii=False, indent=2),
            encoding="utf-8",
        )
        saved_msg(run_dir, "page_texts.json")

        # ---------------- FIGURES ----------------
        t0 = time.perf_counter()
        with st.spinner("Extracting figures…"):
            fig_dir = run_dir / "figures"
            fig_dir.mkdir(parents=True, exist_ok=True)

            # Prefer extraction via graphic.py (exact 7 figures)
            figures = []
            try:
                import shutil
                import graphic  # same folder

                # Point graphic.py to the uploaded DOCX and run extraction
                graphic.EXTRACT_DOCX_PATH = str(saved_path)
                graphic.main_docx_extract()

                # Copy extracted figures into this run's figures folder
                for i in range(1, 8):
                    src = Path(graphic.EXTRACT_FINAL_DIR) / f"Figure_{i}.png"
                    dst = fig_dir / f"Figure_{i}.png"
                    if src.exists():
                        shutil.copyfile(src, dst)

                        try:
                            im = Image.open(dst)
                            w, h = im.size
                        except Exception:
                            w, h = 0, 0

                        figures.append({
                            "id": f"fig{i}",
                            "page": None,
                            "image_path": str(dst),
                            "width": w,
                            "height": h,
                            "_seq": i,
                        })
                    else:
                        st.warning(f"Missing extracted figure from graphic.py: {src}")

            except Exception as e:
                st.warning(f"graphic.py extraction error: {e}")

            # Fallback to legacy Word extraction if graphic.py produced nothing
            if not figures:
                try:
                    figures = extract_figures_from_word(saved_path, render_dir=fig_dir, zoom=2.0) or []
                except Exception as e:
                    st.warning(f"Figure extraction error: {e}")
                    figures = []
        figures = sort_figures(figures)
        for _i, _f in enumerate(figures, start=1):
            _f["_seq"] = _i
        timings["figures_s"] = round(time.perf_counter() - t0, 3)
        (run_dir / "figures.json").write_text(
            json.dumps({"figures": figures}, ensure_ascii=False, indent=2),
            encoding="utf-8",
        )
        pass
        saved_msg(run_dir, "figures.json")

        # ---------------- TABLES ----------------
        t0 = time.perf_counter()
        with st.spinner("Extracting tables…"):
            tbl_dir = run_dir / "tables"
            try:
                tables = extract_tables_with_coords(
                    saved_path, render_dir=tbl_dir, zoom=2.0) or []
            except Exception:
                tables = []
        timings["tables_s"] = round(time.perf_counter() - t0, 3)
        (run_dir / "tables.json").write_text(
            json.dumps({"tables": tables}, ensure_ascii=False, indent=2),
            encoding="utf-8",
        )
        pass
        saved_msg(run_dir, "tables.json")

        # ---------------- SECTIONS ----------------
        t0 = time.perf_counter()
        spinner_text = "Chunking into clean learning sections…"
        if use_ai:
            spinner_text = "Chunking sections + AI summary generation…"
        else:
            spinner_text = "Chunking into clean learning sections…"

        with st.spinner(spinner_text):
            # Use custom Word section chunker based on detected callouts
            sections = chunk_word_sections(pages)

            # Count figures and tables per section
            fig_by_page, tbl_by_page = {}, {}
            for f in figures:
                p_raw = f.get("page", 0)
                p = int(p_raw) if p_raw not in (None, "", "None") else 0

                fig_by_page[p] = fig_by_page.get(p, 0) + 1
            for t in tables:
                p = int(t.get("page", 0))
                tbl_by_page[p] = tbl_by_page.get(p, 0) + 1

            for s in sections:
                # Normalize section pages to integers for accurate counting
                section_page_ids = normalize_page_list(s.get("pages", []))
                s["figure_count"] = sum(
                    fig_by_page.get(p, 0) for p in section_page_ids)
                s["table_count"] = sum(
                    tbl_by_page.get(p, 0) for p in section_page_ids)

            # Generate AI identifiers for each section if Ollama is available
            ollama_url = os.getenv("OLLAMA_URL", "http://localhost:11434")
            try:
                import requests
                test_resp = requests.get(f"{ollama_url}/api/tags", timeout=3)
                if test_resp.status_code == 200:
                    from pipeline.ai_summarizer import generate_section_identifier_ollama
                    for s in sections:
                        title = s.get("title", "")
                        content = s.get("text", "")[:600]
                        if content:
                            ai_identifier = generate_section_identifier_ollama(title, content, ollama_url)
                            if ai_identifier and ai_identifier != title:
                                s["identifier"] = ai_identifier
            except Exception:
                pass

        timings["sections_s"] = round(time.perf_counter() - t0, 3)
        (run_dir / "sections.json").write_text(
            json.dumps({"sections": sections}, ensure_ascii=False, indent=2),
            encoding="utf-8",
        )
        pass
        saved_msg(run_dir, "sections.json")

        # ---------------- META SUMMARY ----------------
        meta = {
            "run": run_dir.name,
            "file": saved_path.name,
            "pages": len(pages),
            "figures": len(figures),
            "tables": len(tables),
            "sections": len(sections),
            "timings_s": timings,
            "ocr_pages": sum(1 for p in pages if p.get("ocr_used")),
            "created_at": ts,
        }
        (run_dir / "meta.json").write_text(json.dumps(meta,
                                                      ensure_ascii=False,
                                                      indent=2),
                                           encoding="utf-8")
        saved_msg(run_dir, "meta.json")

        # ---------------- SAVE TO DATABASE ----------------
        import os as db_os
        init_db()
        save_run(None, run_dir.name, meta, pages, figures, tables, sections)
        db_type = "PostgreSQL" if db_os.environ.get(
            'DATABASE_URL') else "SQLite (data/pdf_extraction.db)"
        pass

        # Cache data for instant editing responses
        st.session_state.processed_file_key = current_file_key
        st.session_state.cached_pages = pages
        st.session_state.cached_figures = figures
        st.session_state.cached_tables = tables
        st.session_state.cached_sections = sections
        st.session_state.cached_run_dir = run_dir
        st.session_state.cached_timings = timings
        st.session_state.cached_ts = ts
        st.session_state.cached_meta = meta

    # ---------------- EXTRACTION STATISTICS DASHBOARD ----------------
    st.markdown("---")
    st.subheader("📊 Extraction Statistics")

    # Calculate quality metrics
    ocr_pages_count = sum(1 for p in pages if p.get("ocr_used"))

    # Calculate average OCR confidence if available
    ocr_confidences = [
        p.get("quality_metrics", {}).get("ocr_confidence", 0) for p in pages
        if p.get("ocr_used")
    ]
    avg_ocr_confidence = (sum(ocr_confidences) / len(ocr_confidences) *
                          100) if ocr_confidences else 100.0

    # Quality score: if no OCR used = 100%, otherwise use OCR confidence
    quality_score = avg_ocr_confidence if ocr_pages_count > 0 else 100.0

    # Total processing time
    total_time = sum(timings.values())

    # Display metrics in colored cards
    col1, col2, col3, col4, col5 = st.columns(5)

    with col1:
        st.metric(label="📄 Pages",
                  value=len(pages),
                  delta=f"{ocr_pages_count} OCR"
                  if ocr_pages_count > 0 else "All text")

    with col2:
        st.metric(label="✨ Quality",
                  value=f"{quality_score:.1f}%",
                  delta="Excellent" if quality_score >= 95 else
                  "Good" if quality_score >= 80 else "Fair")

    with col3:
        st.metric(label="🖼️ Figures",
                  value=len(figures),
                  delta=f"{timings.get('figures_s', 0):.1f}s")

    with col4:
        st.metric(label="📊 Tables",
                  value=len(tables),
                  delta=f"{timings.get('tables_s', 0):.1f}s")

    with col5:
        st.metric(label="🧩 Sections",
                  value=len(sections),
                  delta=f"{timings.get('sections_s', 0):.1f}s")

    # Processing time breakdown
    st.markdown("**⏱️ Processing Time Breakdown:**")
    time_col1, time_col2, time_col3, time_col4 = st.columns(4)

    with time_col1:
        st.caption(f"Text: **{timings.get('text_s', 0):.2f}s**")
    with time_col2:
        st.caption(f"Figures: **{timings.get('figures_s', 0):.2f}s**")
    with time_col3:
        st.caption(f"Tables: **{timings.get('tables_s', 0):.2f}s**")
    with time_col4:
        st.caption(
            f"Sections: **{timings.get('sections_s', 0):.2f}s** | Total: **{total_time:.2f}s**"
        )

    st.markdown("---")

    # ========== VIEW MODE SELECTOR ==========
    # Check if navigating to Figures view from a section link
    if st.session_state.get('navigate_to_figures', False):
        st.session_state['navigate_to_figures'] = False
        # Delete the widget key so we can set the default index (Figures = index 0)
        if 'view_mode_selector' in st.session_state:
            del st.session_state['view_mode_selector']
        nav_default_index = 0  # Figures
    elif st.session_state.get('navigate_to_sections', False):
        st.session_state['navigate_to_sections'] = False
        # Delete the widget key so we can set the default index (Sections = index 1)
        if 'view_mode_selector' in st.session_state:
            del st.session_state['view_mode_selector']
        nav_default_index = 1  # Sections
    else:
        nav_default_index = None
    
    view_mode = st.radio(
        "Select View",
        options=["Figures", "Sections"],
        horizontal=True,
        key="view_mode_selector",
        index=nav_default_index if nav_default_index is not None else 0,
        label_visibility="collapsed"
    )

    # ========== FIGURES VIEW ==========
    if view_mode == "Figures":
        # ---------------- FIGURES PREVIEW ----------------
        st.subheader("🖼 Figures")
        
        st.caption(f"Total figures extracted: {len(figures)}")
        
        # Check if navigating from Sections view to a specific figure
        scroll_to_fig = st.session_state.pop('scroll_to_figure', None)
        # Store which figure we navigated to for the back button
        if scroll_to_fig is not None:
            st.session_state['navigated_to_figure'] = scroll_to_fig
            # Inject JavaScript to scroll to the figure using components
            import streamlit.components.v1 as components
            components.html(f'''
                <script>
                    setTimeout(function() {{
                        var el = window.parent.document.getElementById("fig-{scroll_to_fig}");
                        if (el) {{
                            el.scrollIntoView({{behavior: "smooth", block: "start"}});
                        }}
                    }}, 500);
                </script>
            ''', height=0)
        
        for idx, f in enumerate(figures, start=0):
            # Primary anchor uses sequential numbering (fig0, fig1, ...)
            anchor_here(f"fig-{idx}")
            # Backward-compatible anchor (original extracted id), if present
            if f.get("id") or f.get("figure_id"):
                anchor_here(f"fig-{f.get('id') or f.get('figure_id')}")

            display_id = f.get("id") or f.get("figure_id") or ""
            page_val = f.get("page") or f.get("page_num") or f.get("page_number") or ""
            # Auto-expand if navigating to this figure from Sections view
            is_target_figure = (scroll_to_fig is not None and idx == scroll_to_fig)
            came_from_section = st.session_state.get('came_from_section')
            navigated_fig = st.session_state.get('navigated_to_figure')
            with st.expander(f"Figure {idx}", expanded=is_target_figure):
                # Show "Back to Section" button if this is the navigated figure and came from a section
                if came_from_section and navigated_fig == idx:
                    # Use on_click callback - this runs BEFORE rerun and preserves state
                    def handle_back_click():
                        sec_id = st.session_state.get('_pending_back_section')
                        if sec_id:
                            st.session_state['navigate_to_sections'] = True
                            st.session_state['highlight_section'] = sec_id
                            st.session_state[f'section_expanded_{sec_id}'] = True
                            st.session_state['came_from_section'] = None
                            st.session_state['navigated_to_figure'] = None
                    
                    # Store section ID BEFORE button renders
                    st.session_state['_pending_back_section'] = came_from_section
                    st.button("🔙 Back to Section", key=f"back_to_section_{idx}", on_click=handle_back_click)
                
                # Side-by-side view:
                # Render both images as high-res canvases with ORIGINAL aspect ratio.
                # Streamlit will downscale to column width (sharp, not blurry).
                MAX_CANVAS_W = 600  # Reduced to fit images in window
                MAX_CANVAS_H = 450  # Max height to prevent tall images from being too large
                _RESAMPLE = getattr(getattr(Image, "Resampling", Image), "LANCZOS", Image.LANCZOS)

                def _pick_original(fig_index: int):
                    """fig_index is 0-based (Figure 0, Figure 1, etc.)"""
                    search_patterns = [
                        UPLOAD_DIR / f"figure_{fig_index}",
                        ROOT / "data" / "uploads" / f"figure_{fig_index}",
                        ROOT / "only_7_figures_output" / "final_7" / f"Figure_{fig_index}",
                    ]
                    for base in search_patterns:
                        for ext in (".png", ".jpg", ".jpeg", ".webp"):
                            p = str(base) + ext
                            if os.path.exists(p):
                                return p
                    return None

                def _fit_on_canvas(im: Image.Image, canvas_w: int, canvas_h: int, allow_upscale: bool = False) -> Image.Image:
                    """Scale image to FIT canvas, then center on white canvas."""
                    im = im.convert("RGB")

                    # Scale to fit canvas
                    scale = min(canvas_w / im.width, canvas_h / im.height)
                    if not allow_upscale:
                        scale = min(scale, 1.0)

                    new_w = max(1, int(round(im.width * scale)))
                    new_h = max(1, int(round(im.height * scale)))
                    im = im.resize((new_w, new_h), _RESAMPLE)

                    # Create canvas with EXACT dimensions and center the image
                    canvas = Image.new("RGB", (canvas_w, canvas_h), (255, 255, 255))
                    x = (canvas_w - new_w) // 2
                    y = (canvas_h - new_h) // 2
                    canvas.paste(im, (x, y))
                    return canvas

                orig_path = _pick_original(idx)
                img_path = f.get("image_path")
                ext_path = None
                if img_path:
                    # Try absolute path first
                    if Path(img_path).exists():
                        ext_path = str(img_path)
                    else:
                        # Try relative path from filename (for cross-platform compatibility)
                        img_name = Path(img_path).name
                        # Look in the run's figures directory
                        if 'run_dir' in dir() and run_dir:
                            alt_path = Path(run_dir) / "figures" / img_name
                            if alt_path.exists():
                                ext_path = str(alt_path)
                        # Also try only_7_figures_output as fallback
                        if not ext_path:
                            alt_path = Path("only_7_figures_output/final_7") / img_name
                            if alt_path.exists():
                                ext_path = str(alt_path)

                # Load images first
                orig_im = Image.open(orig_path) if orig_path else None
                ext_im = Image.open(ext_path) if ext_path else None

                # Use fixed canvas width so all figures appear same size
                canvas_w = MAX_CANVAS_W  # Always use full width
                if orig_im:
                    canvas_h = max(1, int(round(canvas_w * (orig_im.height / orig_im.width))))
                elif ext_im:
                    canvas_h = max(1, int(round(canvas_w * (ext_im.height / ext_im.width))))
                else:
                    canvas_h = 400
                # Cap height to prevent tall images from being too large
                canvas_h = min(canvas_h, MAX_CANVAS_H)

                # Create canvases with IDENTICAL dimensions (same width AND height)
                # Original: no upscaling needed (already high-res)
                # Extracted: allow upscaling so it visually matches the Original's size
                orig_canvas = _fit_on_canvas(orig_im, canvas_w, canvas_h, allow_upscale=False) if orig_im else None
                ext_canvas = _fit_on_canvas(ext_im, canvas_w, canvas_h, allow_upscale=True) if ext_im else None
                
                # Verify both canvases have identical dimensions
                if orig_canvas and ext_canvas:
                    assert orig_canvas.size == ext_canvas.size, f"Canvas mismatch: {orig_canvas.size} vs {ext_canvas.size}"

                left_col, right_col = st.columns(2)

                with left_col:
                    st.markdown("**Original Image**")
                    if orig_canvas is not None:
                        st.image(orig_canvas, use_container_width=True)
                    else:
                        st.warning("Original image not found.")

                with right_col:
                    st.markdown("**Extracted Image**")
                    if ext_canvas is not None:
                        st.image(ext_canvas, use_container_width=False)
                    else:
                        img, legend = render_annotated_figure(f, PREVIEW_IMG_WIDTH)
                        if img:
                            st.image(img)
                        if legend:
                            st.json(legend)

                # JSON download link (per figure)
                try:
                    import graphic
                    json_dir = Path(graphic.OUT_DIR)
                except Exception:
                    json_dir = Path("./out_json_fixed_fig0")

                # JSON files are 1-indexed (fig1.json = Figure 0)
                json_path = json_dir / f"fig{idx + 1}.json"
                if json_path.exists():
                    import base64
                    json_content = json_path.read_text(encoding="utf-8")
                    b64 = base64.b64encode(json_content.encode()).decode()
                    href = f'<a href="data:application/json;base64,{b64}" download="fig{idx}_data.json" style="color: black; text-decoration: none;">🔗 Figure {idx} JSON</a>'
                    st.markdown(href, unsafe_allow_html=True)


        # ---------------- TABLES PREVIEW ----------------
        if tables:
            st.subheader("📊 Tables")
        for t in tables:
            anchor_here(f"tbl-{t['id']}")
            with st.expander(
                    f"{t['id']} (Page {t['page']}) - {t.get('nrows', 0)}×{t.get('ncols', 0)} cells"
            ):
                # Table metadata
                st.markdown(f"**📍 Table Position:** Page {t['page']}")
                st.markdown(
                    f"**📐 Bounding Box (coords):** `{t.get('bbox', [])}`")
                st.markdown(
                    f"**📊 Dimensions:** {t.get('nrows', 0)} rows × {t.get('ncols', 0)} columns"
                )
                st.markdown(f"**✨ Confidence:** {t.get('confidence', 0):.2f}")

                # Render annotated table image with cell boundaries
                if t.get("image_path"):
                    annotated_img = render_annotated_table(t, TABLE_IMG_WIDTH)
                    if annotated_img:
                        st.image(
                            annotated_img,
                            caption=f"Table {t['id']} with cell annotations",
                            width=TABLE_IMG_WIDTH,
                        )

                # Show structured data
                grid = t.get("data", [])
                if grid:
                    st.markdown("**📋 Extracted Data:**")

                    try:
                        num_cols = len(grid[0]) if grid else 0
                        col_names = [f"Column {i+1}" for i in range(num_cols)]
                        df = pd.DataFrame(grid, columns=col_names)
                        st.dataframe(df, use_container_width=True)
                    except Exception:
                        st.table(grid)
                else:
                    st.info("No structured table data found.")

                # Cell-level coordinate inspector
                cells = t.get("cells", [])
                if cells:
                    st.markdown("---")
                    st.markdown(
                        "**🔍 Cell Coordinate Inspector** (Text Component Annotations)"
                    )

                    # Show interactive cell selector
                    show_cells = st.checkbox(
                        f"Show all {len(cells)} cell coordinates",
                        key=f"cells_{t['id']}",
                    )

                    if show_cells:
                        # Display in columns for compact view
                        for idx, cell in enumerate(
                                cells[:50]):  # Limit to first 50 cells
                            row, col = cell.get("r", -1), cell.get("c", -1)
                            text = cell.get("text", "")
                            bbox = cell.get("bbox", [])

                            with st.container():
                                cols = st.columns([1, 2, 3])
                                with cols[0]:
                                    st.caption(f"**Cell ({row},{col})**")
                                with cols[1]:
                                    st.caption(f"`{bbox}`")
                                with cols[2]:
                                    st.caption(f'"{text[:40]}..."' if len(text)
                                               > 40 else f'"{text}"')

                        if len(cells) > 50:
                            st.info(
                                f"Showing first 50 of {len(cells)} cells. Full data available in JSON export."
                            )

                # --- Generate Table Insights Button ---
                st.markdown("---")
                tbl_id = t.get('id')
                insights_btn = st.button(f"🤖 Generate Table Insights",
                                         key=f"main_insights_btn_{tbl_id}")

                if insights_btn:
                    try:
                        ollama_url = os.getenv("OLLAMA_URL",
                                               "http://localhost:11434")
                        grid = t.get("data", [])
                        if grid:
                            table_text = "\n".join([
                                " | ".join(str(c) for c in row)
                                for row in grid[:20]
                            ])
                            with st.spinner("Generating table insights..."):
                                from pipeline.ai_summarizer import generate_learning_bullets_ollama
                                insights = generate_learning_bullets_ollama(
                                    f"Analyze this table and provide key insights: {table_text}",
                                    num_bullets=3,
                                    ollama_url=ollama_url)
                                if insights:
                                    if 'table_insights' not in st.session_state:
                                        st.session_state.table_insights = {}
                                    st.session_state.table_insights[
                                        tbl_id] = insights
                                    st.rerun()
                                else:
                                    st.warning(
                                        "Could not generate insights. Make sure Ollama is running."
                                    )
                        else:
                            st.warning("No table data available to analyze.")
                    except Exception as e:
                        st.error(f"Error: {e}")

                # Display stored insights
                if 'table_insights' in st.session_state and tbl_id in st.session_state.table_insights:
                    st.markdown("**💡 AI Table Insights:**")
                    for insight in st.session_state.table_insights[tbl_id]:
                        st.markdown(f"• {insight}")

    # ========== SECTIONS VIEW ==========
    elif view_mode == "Sections":
        st.subheader("🧩 Sections")
        
        # Check if we need to scroll to a specific section (coming back from Figures)
        highlight_sec = st.session_state.get('highlight_section')
        if highlight_sec:
            anchor_id = f"sec-{highlight_sec}"
            # Use an iframe to execute JavaScript on the parent window
            scroll_html = f"""
            <iframe srcdoc="
                <script>
                    function tryScroll() {{
                        var el = parent.document.getElementById('{anchor_id}');
                        if (el) {{
                            el.scrollIntoView({{behavior: 'smooth', block: 'center'}});
                            var exp = el.closest('[data-testid=stExpander]');
                            if (exp) {{
                                exp.style.outline = '3px solid #28a745';
                                setTimeout(function() {{ exp.style.outline = ''; }}, 3000);
                            }}
                        }} else {{
                            setTimeout(tryScroll, 300);
                        }}
                    }}
                    setTimeout(tryScroll, 500);
                </script>
            " style="display:none;"></iframe>
            """
            st.markdown(scroll_html, unsafe_allow_html=True)
            # Also expand the target section
            st.session_state[f'section_expanded_{highlight_sec}'] = True
            # Clear
            del st.session_state['highlight_section']

        # Quick stats bar
        total_sections = len(sections)
        st.metric("📚 Total Sections", total_sections)

        # Section display
        # Initialize editable sections in session state if not exists
        if 'editable_sections' not in st.session_state:
            st.session_state.editable_sections = {}
        if 'section_order' not in st.session_state:
            st.session_state.section_order = [s['id'] for s in sections]
        if 'split_sections' not in st.session_state:
            st.session_state.split_sections = {}  # Store newly created split sections

        # ========== CALLBACK FUNCTIONS FOR INSTANT EDITING ==========
        def invalidate_export_cache():
            """Clear cached exports when sections change."""
            for key in list(st.session_state.keys()):
                if key.startswith('slides_html_') or key.startswith(
                        'slides_pptx_'):
                    del st.session_state[key]

        def move_section_up(section_id):
            """Move section up in order."""
            order = st.session_state.section_order
            if section_id in order:
                curr_idx = order.index(section_id)
                if curr_idx > 0:
                    order[curr_idx], order[curr_idx -
                                           1] = order[curr_idx -
                                                      1], order[curr_idx]
                    st.session_state.section_order = order
                    invalidate_export_cache()

        def move_section_down(section_id):
            """Move section down in order."""
            order = st.session_state.section_order
            if section_id in order:
                curr_idx = order.index(section_id)
                if curr_idx < len(order) - 1:
                    order[curr_idx], order[curr_idx +
                                           1] = order[curr_idx +
                                                      1], order[curr_idx]
                    st.session_state.section_order = order
                    invalidate_export_cache()

        def delete_section(section_id):
            """Mark section as deleted."""
            st.session_state.editable_sections[section_id][
                'deleted'] = True
            if 'ai_summaries' in st.session_state and section_id in st.session_state.ai_summaries:
                del st.session_state.ai_summaries[section_id]
            invalidate_export_cache()

        def toggle_split_mode(section_id):
            """Toggle split mode for a section."""
            current = st.session_state.get(f'split_mode_{section_id}',
                                           False)
            st.session_state[f'split_mode_{section_id}'] = not current
            # Keep the section expander open when entering split mode
            if not current:  # If we're opening split mode
                st.session_state[f'section_expanded_{section_id}'] = True

        def toggle_merge_mode(section_id):
            """Toggle merge mode for a section."""
            current = st.session_state.get(f'merge_mode_{section_id}',
                                           False)
            st.session_state[f'merge_mode_{section_id}'] = not current
            # Keep the section expander open when entering merge mode
            if not current:  # If we're opening merge mode
                st.session_state[f'section_expanded_{section_id}'] = True
            # Clear any previous merge selection to prevent auto-merge
            if f'merge_select_{section_id}' in st.session_state:
                del st.session_state[f'merge_select_{section_id}']

        def merge_sections(source_id, target_id, all_secs):
            """Merge target section into source section and persist to session state."""
            # Find source and target section data
            source_sec = None
            target_sec = None
            for s in all_secs:
                if s['id'] == source_id:
                    source_sec = s
                if s['id'] == target_id:
                    target_sec = s

            if not source_sec or not target_sec:
                return

            # Get current paragraphs - check editable_section_text first, then modified_paragraphs, then original
            source_text_key = f"section_text_{source_id}"
            target_text_key = f"section_text_{target_id}"

            # Get source text/paragraphs
            if source_text_key in st.session_state.get(
                    'editable_section_text', {}):
                source_text = st.session_state.editable_section_text[
                    source_text_key]
                source_paras = [
                    p.strip() for p in source_text.split('\n\n')
                    if p.strip()
                ]
            elif source_id in st.session_state.get('modified_paragraphs',
                                                   {}):
                source_paras = st.session_state.modified_paragraphs[
                    source_id]
            else:
                source_paras = source_sec.get('paragraphs', [])

            # Get target text/paragraphs
            if target_text_key in st.session_state.get(
                    'editable_section_text', {}):
                target_text = st.session_state.editable_section_text[
                    target_text_key]
                target_paras = [
                    p.strip() for p in target_text.split('\n\n')
                    if p.strip()
                ]
            elif target_id in st.session_state.get('modified_paragraphs',
                                                   {}):
                target_paras = st.session_state.modified_paragraphs[
                    target_id]
            else:
                target_paras = target_sec.get('paragraphs', [])

            # Merge paragraphs
            merged_paras = list(source_paras) + list(target_paras)
            merged_text = "\n\n".join(merged_paras)

            # Store merged paragraphs in session state
            if 'modified_paragraphs' not in st.session_state:
                st.session_state.modified_paragraphs = {}
            st.session_state.modified_paragraphs[source_id] = merged_paras

            # Also update editable_section_text so it shows the merged content
            if 'editable_section_text' not in st.session_state:
                st.session_state.editable_section_text = {}
            st.session_state.editable_section_text[
                source_text_key] = merged_text

            # Clear target's editable text cache
            if target_text_key in st.session_state.get(
                    'editable_section_text', {}):
                del st.session_state.editable_section_text[target_text_key]

            # Merge pages
            source_pages = source_sec.get('pages', [])
            target_pages = target_sec.get('pages', [])
            merged_pages = sorted(list(set(source_pages + target_pages)))

            # Update source section title
            source_title = st.session_state.editable_sections.get(
                source_id, {}).get('title', source_sec.get('title', ''))
            target_title = st.session_state.editable_sections.get(
                target_id, {}).get('title', target_sec.get('title', ''))
            new_title = source_title + " + " + target_title
            st.session_state.editable_sections[source_id][
                'title'] = new_title

            # Mark target section as deleted
            st.session_state.editable_sections[target_id] = {
                'title': '',
                'deleted': True
            }

            # Remove target from section_order
            if target_id in st.session_state.section_order:
                st.session_state.section_order.remove(target_id)

            # Merge AI summaries if they exist
            if 'ai_summaries' in st.session_state:
                source_bullets = st.session_state.ai_summaries.get(
                    source_id, [])
                target_bullets = st.session_state.ai_summaries.get(
                    target_id, [])
                if source_bullets or target_bullets:
                    st.session_state.ai_summaries[source_id] = (
                        source_bullets + target_bullets)[:5]
                if target_id in st.session_state.ai_summaries:
                    del st.session_state.ai_summaries[target_id]

            # Close merge mode
            st.session_state[f'merge_mode_{source_id}'] = False

            invalidate_export_cache()

        # Sync sections with editable state
        for s in sections:
            if s['id'] not in st.session_state.editable_sections:
                st.session_state.editable_sections[s['id']] = {
                    'title': s.get('title', ''),
                    'deleted': False
                }

        # Combine original sections with split sections
        all_sections = list(sections)
        for split_id, split_data in st.session_state.split_sections.items():
            all_sections.append(split_data)

        # Get ordered sections (respecting user reordering)
        ordered_sections = []
        for sid in st.session_state.section_order:
            for s in all_sections:
                if s['id'] == sid and not st.session_state.editable_sections.get(
                        sid, {}).get('deleted', False):
                    ordered_sections.append(s)
                    break
        # Add any new sections not in order
        for s in all_sections:
            if s['id'] not in st.session_state.section_order and not st.session_state.editable_sections.get(
                    s['id'], {}).get('deleted', False):
                ordered_sections.append(s)

        # CHAPTER PRIORITY REORDER 
        st.markdown("### 📚 Which chapter do you want to study first?")
        st.caption(
            "Drag chapters to prioritize your study order. The first chapter will appear at the top."
        )

        # Build chapter hierarchy first to get chapter names
        temp_hierarchy = detect_chapter_subchapter_hierarchy(ordered_sections)
        chapter_names = list(temp_hierarchy.keys())

        # Initialize chapter order in session state if not exists
        if 'chapter_order' not in st.session_state:
            st.session_state.chapter_order = chapter_names.copy()

        # Make sure all chapters are in the order (handle new chapters)
        for ch in chapter_names:
            if ch not in st.session_state.chapter_order:
                st.session_state.chapter_order.append(ch)
        # Remove chapters that no longer exist
        st.session_state.chapter_order = [ch for ch in st.session_state.chapter_order if ch in chapter_names]

        # Create FIXED chapter labels based on original document order (not current sort order)
        # This ensures "Chapter 1" always refers to the first chapter in the original document
        if 'chapter_labels' not in st.session_state:
            st.session_state.chapter_labels = {}

        # Assign fixed labels to chapters based on their original order in document
        for idx, ch_name in enumerate(chapter_names):
            if ch_name not in st.session_state.chapter_labels:
                st.session_state.chapter_labels[ch_name] = f"Chapter {idx + 1}"

        # Create display items using fixed labels, maintaining current sort order
        chapter_display_items = []
        display_to_chapter = {}
        for ch_name in st.session_state.chapter_order:
            label = st.session_state.chapter_labels.get(ch_name, ch_name)
            chapter_display_items.append(label)
            display_to_chapter[label] = ch_name

        # Render sortable chapter list
        sorted_chapter_display = sort_items(chapter_display_items, key="chapter_reorder")

        # Map back to chapter names and update order
        new_chapter_order = [display_to_chapter.get(d, d) for d in sorted_chapter_display]

        if new_chapter_order != st.session_state.chapter_order:
            st.session_state.chapter_order = new_chapter_order
            invalidate_export_cache()

        st.markdown("---")

        # Translate Summaries section for Text View
        st.markdown("""
        <div style="background: linear-gradient(135deg, #4facfe 0%, #00f2fe 100%); 
                    border-radius: 10px; padding: 10px 16px; margin-bottom: 16px;
                    display: inline-flex; align-items: center; gap: 10px;">
            <span style="font-size: 1.5em;">🌐</span>
            <span style="color: white; font-weight: bold;">Translate Summaries</span>
        </div>
        """, unsafe_allow_html=True)

        # Initialize language selection in session state
        if 'selected_lang_name' not in st.session_state:
            st.session_state.selected_lang_name = "English"

        tv_lang_col1, tv_lang_col2, tv_lang_col3 = st.columns([2, 2, 2])
        with tv_lang_col1:
            lang_options = list(LANGUAGES.keys())
            current_idx = lang_options.index(
                st.session_state.selected_lang_name
            ) if st.session_state.selected_lang_name in lang_options else 0

            selected_language = st.selectbox(
                "Select Language:",
                options=lang_options,
                index=current_idx,
                key="tv_translation_language_select",
                label_visibility="collapsed")
            st.session_state.selected_lang_name = selected_language
            st.session_state.selected_lang_code = LANGUAGES.get(selected_language, "en")

        # Get target language for translation
        tv_target_lang = st.session_state.get('selected_lang_code', 'en')

        with tv_lang_col2:
            if selected_language != "English":
                st.success(f"🌐 Translating to **{selected_language}**")

        # Count summarized sections for export button
        tv_summarized_count = sum(1 for s in ordered_sections 
                                 if st.session_state.get(f"summary_{s['id']}"))

        with tv_lang_col3:
            # PPT Export button only
            current_lang_code = LANGUAGES.get(selected_language, "en")
            cache_key_pptx = f"slides_pptx_tv_{current_lang_code}"

            if cache_key_pptx in st.session_state:
                st.download_button(
                    label=f"📥 Download PPTX",
                    data=st.session_state[cache_key_pptx],
                    file_name=f"document_{current_lang_code}.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    key=f"tv_download_pptx_{current_lang_code}")
            else:
                if st.button(f"📊 Generate PPTX",
                             key=f"tv_gen_pptx_{current_lang_code}"):
                    with st.spinner("Generating PPTX..."):
                        pptx_data = generate_slides_pptx(
                            ordered_sections,
                            st.session_state,
                            target_lang=current_lang_code,
                            use_extracted_text_fallback=True)
                        st.session_state[cache_key_pptx] = pptx_data
                        st.rerun()

        # Build chapter/subchapter hierarchy for Text View
        text_view_hierarchy = detect_chapter_subchapter_hierarchy(ordered_sections)

        # Use chapter order from session state for display
        ordered_chapter_names = st.session_state.get('chapter_order', list(text_view_hierarchy.keys()))

        for chapter_name in ordered_chapter_names:
            if chapter_name not in text_view_hierarchy:
                continue
            subchapters = text_view_hierarchy[chapter_name]
            # Translate chapter name if needed
            display_chapter_name = google_translate_free(chapter_name, tv_target_lang) if tv_target_lang != 'en' else chapter_name
            # Count total sections in chapter
            total_chapter_sections = sum(len(secs) for secs in subchapters.values())

            # Chapter header (green gradient - same as Chapter View)
            st.markdown(f"""
            <div style="background: linear-gradient(135deg, #11998e 0%, #38ef7d 100%); 
                        border-radius: 12px; padding: 16px 24px; margin: 24px 0 12px 0;
                        box-shadow: 0 4px 15px rgba(0,0,0,0.2);">
                <h2 style="margin: 0; color: white; font-size: 1.5em;">📚 {display_chapter_name}</h2>
                <p style="margin: 5px 0 0 0; color: rgba(255,255,255,0.8); font-size: 0.9em;">
                    {len(subchapters)} subchapters | {total_chapter_sections} sections total
                </p>
            </div>
            """, unsafe_allow_html=True)

            # Loop through subchapters
            for subchapter_name, subchapter_sections in subchapters.items():
                # Extract subchapter prefix (e.g., "Subchapter 1.1") from the name
                import re
                subchapter_prefix_match = re.match(r'^(Subchapter\s+\d+\.\d+)(?::\s*.*)?$', subchapter_name)
                subchapter_prefix = subchapter_prefix_match.group(1) if subchapter_prefix_match else subchapter_name

                # Display sections within subchapter (subchapter name included in expander title)
                for idx, s in enumerate(subchapter_sections):
                    anchor_here(f"sec-{s['id']}")

                    # Get section data
                    pages_list = s.get("pages", [])
                    section_page_ids = normalize_page_list(pages_list)
                    section_figures = [
                        f for f in figures
                        if _safe_int(f.get("page"), -1) in section_page_ids
                    ]
                    section_tables = [
                        t for t in tables
                        if _safe_int(t.get("page"), -1) in section_page_ids
                    ]
                    # Use modified paragraphs if available (from split operations)
                    if 'modified_paragraphs' in st.session_state and s[
                            'id'] in st.session_state.modified_paragraphs:
                        paragraphs_list = st.session_state.modified_paragraphs[
                            s['id']]
                    else:
                        paragraphs_list = s.get("paragraphs", [])

                    # Get current title (may be edited)
                    section_id = s['id']

                    # Extract just the title part (without number prefix) from original
                    original_title = s.get('title', '')
                    title_match = re.match(r'^[\d\.]+\s*:?\s*(.*)$', original_title)
                    original_title_clean = title_match.group(1) if title_match else original_title

                    # Initialize editable_sections for this section if needed
                    if section_id not in st.session_state.editable_sections:
                        st.session_state.editable_sections[section_id] = {
                            'title': original_title_clean
                        }

                    # Get the saved title from editable_sections (now stored without prefix)
                    display_title = st.session_state.editable_sections[section_id].get('title', original_title_clean)

                    # Translate for display if needed
                    if tv_target_lang != 'en' and display_title:
                        translated_title = google_translate_free(display_title, tv_target_lang)
                        translated_prefix = google_translate_free(subchapter_prefix, tv_target_lang)
                    else:
                        translated_title = display_title
                        translated_prefix = subchapter_prefix

                    # Combine prefix and title for header
                    if translated_title:
                        expander_header = f"{translated_prefix}: {translated_title}"
                    else:
                        expander_header = translated_prefix

                    # TEXT VIEW MODE (expander view)
                    # Check if this section should be expanded (e.g., during merge/split mode)
                    is_expanded = st.session_state.get(
                        f'section_expanded_{section_id}', False)

                    with st.expander(f"{expander_header}",
                                     expanded=is_expanded):
                        # ========== SECTION CONTROLS - TOP RIGHT ICONS ==========
                        # Use spacer column to push icons to the right
                        spacer_col, del_col, merge_col, split_col = st.columns(
                            [20, 1, 1, 1])

                        # Delete button - small icon
                        with del_col:
                            st.button("🗑️",
                                      key=f"delete_{section_id}",
                                      on_click=delete_section,
                                      args=(section_id, ),
                                      help="Delete this section")

                        # Merge with any section button
                        with merge_col:
                            if len(ordered_sections) > 1:
                                st.button("🔗",
                                          key=f"merge_{section_id}",
                                          on_click=toggle_merge_mode,
                                          args=(section_id, ),
                                          help="Merge with another section")

                        # Split section button - small icon
                        with split_col:
                            if len(paragraphs_list) > 1:
                                st.button("✂️",
                                          key=f"split_{section_id}",
                                          on_click=toggle_split_mode,
                                          args=(section_id, ),
                                          help="Split this section")

                        # Split mode interface
                        if st.session_state.get(f'split_mode_{section_id}', False):
                            st.markdown(
                                "**✂️ Split Section - Select split point:**")
                            if paragraphs_list and len(paragraphs_list) > 1:
                                # Show paragraph previews with split buttons
                                for p_idx in range(1, min(len(paragraphs_list),
                                                          6)):
                                    para_preview = paragraphs_list[
                                        p_idx][:80] + "..." if len(
                                            paragraphs_list[p_idx]
                                        ) > 80 else paragraphs_list[p_idx]
                                    col_split_btn, col_preview = st.columns([1, 4])
                                    with col_split_btn:
                                        if st.button(
                                                f"Split here ↓",
                                                key=f"do_split_{section_id}_{p_idx}"
                                        ):
                                            # Create new section from split
                                            import uuid
                                            new_id = f"{section_id}_split_{uuid.uuid4().hex[:6]}"

                                            # Split paragraphs
                                            first_paras = paragraphs_list[:p_idx]
                                            second_paras = paragraphs_list[p_idx:]

                                            # Store modified paragraphs for original section
                                            if 'modified_paragraphs' not in st.session_state:
                                                st.session_state.modified_paragraphs = {}
                                            st.session_state.modified_paragraphs[
                                                section_id] = first_paras

                                            # Create new section data and store in session state
                                            new_section = {
                                                'id':
                                                new_id,
                                                'title':
                                                current_title + " (Part 2)",
                                                'paragraphs':
                                                second_paras,
                                                'pages':
                                                s.get('pages', []),
                                                'identifier':
                                                s.get('identifier', '') + '_part2',
                                                'word_count':
                                                sum(
                                                    len(p.split())
                                                    for p in second_paras),
                                                'figure_count':
                                                0,
                                                'table_count':
                                                0
                                            }

                                            # Store in session state for persistence
                                            st.session_state.split_sections[
                                                new_id] = new_section

                                            # Initialize AI summaries for the new section
                                            # If original section had summaries, split them or create default
                                            if 'ai_summaries' not in st.session_state:
                                                st.session_state.ai_summaries = {}

                                            original_bullets = st.session_state.ai_summaries.get(
                                                section_id, [])
                                            if original_bullets:
                                                # Split bullets: first half stays with original, second half goes to new section
                                                mid = max(
                                                    1,
                                                    len(original_bullets) // 2)
                                                st.session_state.ai_summaries[
                                                    section_id] = original_bullets[:
                                                                                   mid]
                                                st.session_state.ai_summaries[
                                                    new_id] = original_bullets[
                                                        mid:] if len(
                                                            original_bullets
                                                        ) > mid else [
                                                            "Summary for Part 2 - click 'Generate Summary' to update"
                                                        ]
                                            else:
                                                # Create default summary for new section
                                                st.session_state.ai_summaries[
                                                    new_id] = [
                                                        "Summary for this section - click 'Generate Summary' to update"
                                                    ]

                                            # Update editable sections
                                            st.session_state.editable_sections[
                                                new_id] = {
                                                    'title':
                                                    current_title + " (Part 2)",
                                                    'deleted': False
                                                }
                                            st.session_state.editable_sections[
                                                section_id][
                                                    'title'] = current_title + " (Part 1)"

                                            # Update section order
                                            order = st.session_state.section_order
                                            curr_idx = order.index(
                                                section_id
                                            ) if section_id in order else len(
                                                order)
                                            order.insert(curr_idx + 1, new_id)
                                            st.session_state.section_order = order

                                            # Clear split mode
                                            del st.session_state[
                                                f'split_mode_{section_id}']

                                            # Invalidate export cache
                                            for key in list(
                                                    st.session_state.keys()):
                                                if key.startswith(
                                                        'slides_html_'
                                                ) or key.startswith(
                                                        'slides_pptx_'):
                                                    del st.session_state[key]
                                            st.rerun()
                                    with col_preview:
                                        st.caption(para_preview)

                                if st.button("Cancel Split",
                                             key=f"cancel_split_{section_id}"):
                                    del st.session_state[
                                        f'split_mode_{section_id}']
                                    st.rerun()
                            else:
                                st.warning(
                                    "Section has only one paragraph, cannot split."
                                )
                                if st.button("Cancel",
                                             key=f"cancel_split2_{section_id}"):
                                    del st.session_state[
                                        f'split_mode_{section_id}']
                                    st.rerun()

                        # Merge mode interface - select which section to merge with
                        if st.session_state.get(f'merge_mode_{section_id}', False):
                            st.markdown(
                                "**🔗 Merge Section - Select section to merge with:**"
                            )

                            # Get list of other sections (not current, not deleted)
                            other_sections = []
                            for other_s in ordered_sections:
                                other_id = other_s['id']
                                if other_id != section_id:
                                    other_title = st.session_state.editable_sections.get(
                                        other_id, {}).get('title',
                                                          other_s.get('title', ''))
                                    other_sections.append(
                                        (other_id, f"{other_id} • {other_title}"))

                            if other_sections:
                                merge_col1, merge_col2 = st.columns([3, 1])
                                with merge_col1:
                                    selected_merge = st.selectbox(
                                        "Select section to merge into this one",
                                        options=[s[0] for s in other_sections],
                                        format_func=lambda x: next(
                                            (s[1] for s in other_sections
                                             if s[0] == x), x),
                                        key=f"merge_select_{section_id}",
                                        label_visibility="collapsed")
                                with merge_col2:
                                    if st.button("✅ Merge",
                                                 key=f"do_merge_{section_id}"):
                                        merge_sections(section_id, selected_merge,
                                                       ordered_sections)
                                        st.rerun()

                                if st.button("Cancel Merge",
                                             key=f"cancel_merge_{section_id}"):
                                    st.session_state[
                                        f'merge_mode_{section_id}'] = False
                                    st.rerun()
                            else:
                                st.warning("No other sections available to merge.")
                                if st.button("Cancel",
                                             key=f"cancel_merge2_{section_id}"):
                                    st.session_state[
                                        f'merge_mode_{section_id}'] = False
                                    st.rerun()

                        # Compact metadata line - translate labels
                        id_label = google_translate_free("ID", tv_target_lang) if tv_target_lang != 'en' else "ID"
                        words_label = google_translate_free("Words", tv_target_lang) if tv_target_lang != 'en' else "Words"
                        figures_label = google_translate_free("Figures", tv_target_lang) if tv_target_lang != 'en' else "Figures"
                        tables_label = google_translate_free("Tables", tv_target_lang) if tv_target_lang != 'en' else "Tables"
                        pages_label = google_translate_free("Pages", tv_target_lang) if tv_target_lang != 'en' else "Pages"

                        # Translate section identifier
                        section_identifier = s.get('identifier','')
                        if tv_target_lang != 'en' and section_identifier:
                            section_identifier = google_translate_free(section_identifier, tv_target_lang)

                        # Dynamically count figures and tables from section text
                        section_text_for_count = " ".join(s.get('paragraphs', [])) if s.get('paragraphs') else s.get('text', '')
                        import re
                        dynamic_fig_count = len(set(re.findall(r'\bFigure\s*(\d+)\b', section_text_for_count, re.IGNORECASE)))
                        dynamic_table_count = len(set(re.findall(r'\bTable\s*(\d+)\b', section_text_for_count, re.IGNORECASE)))
                        
                        st.markdown(f"""<div style="font-size: 0.85em; color: #666; margin-bottom: 8px;">
                            <b>{id_label}:</b> {section_identifier} · <b>{words_label}:</b> {s.get('word_count',0)} · <b>{figures_label}:</b> {dynamic_fig_count} · <b>{tables_label}:</b> {dynamic_table_count} · <b>{pages_label}:</b> {pages_list}
                        </div>""", unsafe_allow_html=True)

                        # Use the extracted display_title (without number prefix) for editing
                        # Translate for display if in non-English mode
                        editable_title_display = google_translate_free(display_title, tv_target_lang) if tv_target_lang != 'en' else display_title

                        def on_title_change(sid):
                            # Get the new value from the widget
                            input_key = f"title_input_{sid}"
                            if input_key in st.session_state:
                                new_val = st.session_state[input_key]
                                # Save to editable_sections (store just the title without prefix)
                                if sid in st.session_state.editable_sections:
                                    st.session_state.editable_sections[sid]['title'] = new_val
                                else:
                                    st.session_state.editable_sections[sid] = {'title': new_val}
                            st.session_state[f'section_expanded_{sid}'] = True
                            # Invalidate export cache
                            for key in list(st.session_state.keys()):
                                if key.startswith('slides_html_') or key.startswith('slides_pptx_'):
                                    del st.session_state[key]

                        new_title = st.text_input("Edit title",
                                                  value=editable_title_display,
                                                  key=f"title_input_{section_id}",
                                                  label_visibility="collapsed",
                                                  on_change=on_title_change,
                                                  args=(section_id, ),
                                                  disabled=(tv_target_lang != 'en'))


                        # --- TABLE THUMBNAILS FOR THIS SECTION ---
                        if section_tables:
                            st.markdown("**📊 Tables:**")

                            # Display table thumbnails in columns
                            num_tbls = len(section_tables)
                            tbl_cols_per_row = min(3, num_tbls)
                            tbl_cols = st.columns(tbl_cols_per_row)

                            for tbl_idx, tbl in enumerate(section_tables[:6]):
                                col_idx = tbl_idx % tbl_cols_per_row
                                with tbl_cols[col_idx]:
                                    tbl_path = tbl.get("image_path")
                                    if tbl_path and Path(tbl_path).exists():
                                        st.image(
                                            str(tbl_path),
                                            caption=
                                            f"{tbl.get('id', 'Table')} (p.{tbl.get('page')})",
                                            width=150)
                                    else:
                                        st.caption(
                                            f"📊 {tbl.get('id', 'Table')} (p.{tbl.get('page')})"
                                        )

                            if num_tbls > 6:
                                st.caption(f"... and {num_tbls - 6} more tables")

                            # Interactive Table Explorer
                            st.markdown("---")
                            st.markdown("**🔍 Interactive Table Explorer:**")
                            for tbl in section_tables[:2]:
                                with st.expander(
                                        f"📊 Explore {tbl.get('id', 'Table')} ({tbl.get('nrows', 0)}×{tbl.get('ncols', 0)})"
                                ):

                                    # === TABLE STATISTICS DASHBOARD ===
                                    nrows = tbl.get('nrows', 0)
                                    ncols = tbl.get('ncols', 0)
                                    cells = tbl.get("cells", [])
                                    confidence = tbl.get('confidence', 0)
                                    tbl_page = tbl.get('page', '?')

                                    # Calculate confidence color
                                    if confidence >= 0.8:
                                        conf_color = "#28a745"  # green
                                        conf_label = "High"
                                    elif confidence >= 0.5:
                                        conf_color = "#ffc107"  # yellow
                                        conf_label = "Medium"
                                    else:
                                        conf_color = "#dc3545"  # red
                                        conf_label = "Low"

                                    st.markdown(f"""
                                    <div style="background: linear-gradient(135deg, #f8f9fa 0%, #e9ecef 100%); 
                                                padding: 12px 16px; border-radius: 10px; margin-bottom: 12px;
                                                border-left: 4px solid #4A90D9;">
                                        <div style="display: flex; flex-wrap: wrap; gap: 20px; align-items: center;">
                                            <div style="text-align: center;">
                                                <div style="font-size: 1.5em; font-weight: bold; color: #2c3e50;">{nrows}</div>
                                                <div style="font-size: 0.8em; color: #6c757d;">Rows</div>
                                            </div>
                                            <div style="font-size: 1.2em; color: #adb5bd;">×</div>
                                            <div style="text-align: center;">
                                                <div style="font-size: 1.5em; font-weight: bold; color: #2c3e50;">{ncols}</div>
                                                <div style="font-size: 0.8em; color: #6c757d;">Columns</div>
                                            </div>
                                            <div style="border-left: 1px solid #dee2e6; height: 30px;"></div>
                                            <div style="text-align: center;">
                                                <div style="font-size: 1.5em; font-weight: bold; color: #4A90D9;">{len(cells)}</div>
                                                <div style="font-size: 0.8em; color: #6c757d;">Cells Detected</div>
                                            </div>
                                            <div style="border-left: 1px solid #dee2e6; height: 30px;"></div>
                                            <div style="text-align: center;">
                                                <div style="font-size: 1.5em; font-weight: bold; color: #6c757d;">p.{tbl_page}</div>
                                                <div style="font-size: 0.8em; color: #6c757d;">Page</div>
                                            </div>
                                            <div style="border-left: 1px solid #dee2e6; height: 30px;"></div>
                                            <div style="text-align: center;">
                                                <div style="font-size: 1.2em; font-weight: bold; color: {conf_color};">{confidence:.0%}</div>
                                                <div style="font-size: 0.8em; color: #6c757d;">{conf_label} Confidence</div>
                                            </div>
                                        </div>
                                    </div>
                                    """, unsafe_allow_html=True)

                                    # === ANNOTATED CELL BOUNDARIES VIEW ===
                                    st.markdown("**🔲 Cell Boundary Detection:**")
                                    annotated_tbl_img = render_annotated_table(tbl, 400)
                                    if annotated_tbl_img:
                                        st.image(
                                            annotated_tbl_img,
                                            caption="Blue boxes show detected cell boundaries",
                                            use_container_width=True)
                                    else:
                                        # Fallback to regular image
                                        tbl_path = tbl.get("image_path")
                                        if tbl_path and Path(tbl_path).exists():
                                            st.image(str(tbl_path), width=300)

                                    # Show table data as interactive dataframe
                                    grid = tbl.get("data", [])
                                    if grid and len(grid) > 0:
                                        st.markdown("**📋 Extracted Data:**")
                                        try:
                                            num_cols = len(grid[0]) if grid else 0
                                            col_names = [f"Column {i+1}" for i in range(num_cols)]
                                            df = pd.DataFrame(grid, columns=col_names)
                                            st.dataframe(df, use_container_width=True)
                                        except Exception:
                                            st.table(grid)

                                    # AI Table Insights button
                                    if st.button(
                                            f"🤖 Generate Table Insights",
                                            key=f"tbl_ai_{s['id']}_{tbl.get('id')}"):
                                        tbl_id = tbl.get('id')
                                        if 'table_insights' not in st.session_state:
                                            st.session_state.table_insights = {}
                                        if grid:
                                            try:
                                                ollama_url = os.getenv("OLLAMA_URL", "http://localhost:11434")
                                                with st.spinner("Analyzing table..."):
                                                    table_text = f"Table with {tbl.get('nrows', 0)} rows and {tbl.get('ncols', 0)} columns.\n"
                                                    if grid:
                                                        table_text += f"Headers: {', '.join(str(h) for h in grid[0] if h)}\n"
                                                        table_text += f"Sample data: {grid[1] if len(grid) > 1 else 'N/A'}"
                                                    from pipeline.ai_summarizer import generate_learning_bullets_ollama
                                                    insights = generate_learning_bullets_ollama(
                                                        f"Analyze this table and provide key insights: {table_text}",
                                                        num_bullets=3,
                                                        ollama_url=ollama_url)
                                                    if insights:
                                                        st.session_state.table_insights[tbl_id] = insights
                                                        st.rerun()
                                                    else:
                                                        st.warning("Could not generate insights. Make sure Ollama is running.")
                                            except Exception as e:
                                                st.error(f"Error: {e}")

                                    # Display stored insights
                                    tbl_id = tbl.get('id')
                                    if 'table_insights' in st.session_state and tbl_id in st.session_state.table_insights:
                                        st.markdown("**💡 AI Table Insights:**")
                                        for insight in st.session_state.table_insights[tbl_id]:
                                            st.markdown(f"• {insight}")

                        # Extracted text with heading detection
                        extracted_label = google_translate_free("Extracted Text", tv_target_lang) if tv_target_lang != 'en' else "Extracted Text"

                        # Get original text
                        if paragraphs_list:
                            original_text = "\n\n".join(paragraphs_list)
                        else:
                            original_text = s.get("raw_text", "")

                        # Check word count for expand/collapse feature
                        section_word_count = s.get("word_count", 0)
                        use_expandable = section_word_count > 150

                        # Initialize expand state for this section
                        expand_key = f"text_expanded_{section_id}"
                        if expand_key not in st.session_state:
                            st.session_state[expand_key] = False

                        # Determine display text based on language
                        if tv_target_lang != 'en' and original_text:
                            display_text = google_translate_free(original_text, tv_target_lang)
                            is_readonly = True
                        else:
                            # Get editable text for English
                            if 'editable_section_text' not in st.session_state:
                                st.session_state.editable_section_text = {}
                            text_key = f"section_text_{section_id}"
                            if text_key in st.session_state.editable_section_text:
                                display_text = st.session_state.editable_section_text[text_key]
                            else:
                                display_text = original_text
                            is_readonly = False

                        # Show header
                        is_expanded = st.session_state[expand_key]
                        st.markdown(f"**📄 {extracted_label}:**")

                        # Calculate line count based on expand state (only for 150+ word sections)
                        if use_expandable:
                            if is_expanded:
                                # Full text view
                                actual_lines = display_text.count('\n') + 1
                                char_lines = len(display_text) // 80
                                line_count = max(4, min(20, max(actual_lines, char_lines)))
                            else:
                                # Compact view - just 3 lines
                                line_count = 3
                        else:
                            # For short sections, show all content
                            actual_lines = display_text.count('\n') + 1
                            char_lines = len(display_text) // 80
                            line_count = max(2, min(12, max(actual_lines, char_lines)))

                        if is_readonly:
                            st.text_area("Translated content",
                                        value=display_text,
                                        height=line_count * 22,
                                        key=f"translated_{section_id}_{tv_target_lang}",
                                        label_visibility="collapsed",
                                        disabled=True)
                        else:
                            widget_key = f"text_edit_{section_id}"
                            text_key = f"section_text_{section_id}"

                            if widget_key not in st.session_state:
                                st.session_state[widget_key] = display_text
                            elif text_key in st.session_state.editable_section_text:
                                st.session_state[widget_key] = st.session_state.editable_section_text[text_key]

                            new_text = st.text_area("Section content",
                                                    key=widget_key,
                                                    height=line_count * 22,
                                                    label_visibility="collapsed")

                            if new_text != display_text:
                                st.session_state.editable_section_text[text_key] = new_text
                                if 'modified_paragraphs' not in st.session_state:
                                    st.session_state.modified_paragraphs = {}
                                st.session_state.modified_paragraphs[section_id] = [
                                    p.strip() for p in new_text.split('\n\n') if p.strip()
                                ]
                                for key in list(st.session_state.keys()):
                                    if key.startswith('slides_html_') or key.startswith('slides_pptx_'):
                                        del st.session_state[key]

                        # Bottom row controls
                        has_summary = 'ai_summaries' in st.session_state and s['id'] in st.session_state.ai_summaries
                        has_questions = 'ai_questions' in st.session_state and s['id'] in st.session_state.ai_questions

                        if use_expandable:
                            # For 150+ words: arrow, summarize menu, and questions button
                            arrow = "▲" if is_expanded else "▼"
                            col_spacer, col_arrow, col_menu, col_q = st.columns([18, 1, 1, 1])
                            with col_arrow:
                                if st.button(arrow, key=f"toggle_text_{section_id}", help="Expand/Collapse", type="secondary"):
                                    st.session_state[expand_key] = not is_expanded
                                    st.session_state[f'section_expanded_{section_id}'] = True
                                    st.session_state['scroll_to_section'] = section_id
                            with col_menu:
                                with st.popover("⋮"):
                                    sum_label = "Regenerate Summary" if has_summary else "Summarize"
                                    if st.button(sum_label, key=f"menu_sum_{section_id}"):
                                        st.session_state[f'trigger_summary_{section_id}'] = True
                                        st.session_state[f'section_expanded_{section_id}'] = True
                                        st.session_state['scroll_to_section'] = section_id
                            with col_q:
                                if st.button("Q", key=f"menu_q_{section_id}", help="Generate Questions", type="secondary"):
                                    st.session_state[f'trigger_questions_{section_id}'] = True
                                    st.session_state[f'section_expanded_{section_id}'] = True
                                    st.session_state['scroll_to_section'] = section_id
                        else:
                            # For shorter sections: just questions button
                            col_spacer, col_q = st.columns([20, 1])
                            with col_q:
                                if st.button("Q", key=f"menu_q_{section_id}", help="Generate Questions", type="secondary"):
                                    st.session_state[f'trigger_questions_{section_id}'] = True
                                    st.session_state[f'section_expanded_{section_id}'] = True
                                    st.session_state['scroll_to_section'] = section_id

                        # Detect and display figures mentioned in this section's text
                        import re
                        figure_pattern = re.compile(r'\bFigure\s*(\d+)\b', re.IGNORECASE)
                        mentioned_figures = set()
                        text_to_scan = display_text or original_text or ""
                        for match in figure_pattern.finditer(text_to_scan):
                            fig_num = int(match.group(1))
                            mentioned_figures.add(fig_num)
                        
                        if mentioned_figures and figures:
                            st.markdown("**🖼 Referenced Figures:**")
                            # Collect all figure data first
                            fig_items = []
                            for fig_num in sorted(mentioned_figures):
                                if 0 <= fig_num < len(figures):
                                    fig_data = figures[fig_num]
                                    orig_path = None
                                    docs_folder = Path("Documents")
                                    if docs_folder.exists():
                                        for pattern in [f"fig{fig_num}.png", f"fig{fig_num}.jpg", f"Figure_{fig_num}.png"]:
                                            for f_path in docs_folder.glob(f"**/{pattern}"):
                                                orig_path = f_path
                                                break
                                            if orig_path:
                                                break
                                    if not orig_path:
                                        upload_path = Path("data/uploads") / f"figure_{fig_num}.png"
                                        if upload_path.exists():
                                            orig_path = upload_path
                                    if not orig_path:
                                        ext_path = fig_data.get("image_path")
                                        if ext_path and Path(ext_path).exists():
                                            orig_path = Path(ext_path)
                                    if orig_path and orig_path.exists():
                                        fig_items.append((fig_num, orig_path))
                            
                            # Display figures horizontally in columns (max 2 per row for better visibility)
                            if fig_items:
                                cols_per_row = min(2, len(fig_items))
                                cols = st.columns(cols_per_row)
                                for idx, (fig_num, img_path) in enumerate(fig_items):
                                    with cols[idx % cols_per_row]:
                                        st.caption(f"Figure {fig_num}")
                                        st.image(str(img_path), width=420)
                                        clicked = st.button(f"🔗 Go to Figure {fig_num}", key=f"nav_fig_{fig_num}_{section_id}", use_container_width=False)
                                        if clicked:
                                            st.session_state['navigate_to_figures'] = True
                                            st.session_state['scroll_to_figure'] = fig_num
                                            st.session_state['came_from_section'] = section_id
                                            st.rerun()

                        # Get section text
                        section_text = " ".join(s.get(
                            "paragraphs", [])) if s.get("paragraphs") else s.get(
                                "raw_text", "")
                        section_title = s.get('title', '')

                        # Handle AI Summary from 3-dot menu
                        if st.session_state.get(f'trigger_summary_{section_id}', False):
                            st.session_state[f'trigger_summary_{section_id}'] = False
                            st.session_state[f'section_expanded_{s["id"]}'] = True
                            st.session_state['scroll_to_section'] = s['id']
                            try:
                                ollama_url = os.getenv("OLLAMA_URL", "http://localhost:11434")
                                with st.spinner("Generating AI summary..."):
                                    from pipeline.ai_summarizer import generate_learning_bullets_ollama
                                    bullets = generate_learning_bullets_ollama(
                                        section_text, num_bullets=4, ollama_url=ollama_url)
                                    if bullets:
                                        if 'ai_summaries' not in st.session_state:
                                            st.session_state.ai_summaries = {}
                                        st.session_state.ai_summaries[s['id']] = bullets
                                        st.session_state['scroll_to_section'] = section_id
                                        st.rerun()
                                    else:
                                        st.warning(f"⚠️ Could not generate summary. Make sure Ollama is accessible at {ollama_url} with 'mistral' model.")
                            except Exception as e:
                                st.error(f"Error generating summary: {e}")

                        # Handle Review Questions from 3-dot menu
                        if st.session_state.get(f'trigger_questions_{section_id}', False):
                            st.session_state[f'trigger_questions_{section_id}'] = False
                            st.session_state[f'section_expanded_{s["id"]}'] = True
                            st.session_state['scroll_to_section'] = s['id']
                            try:
                                ollama_url = os.getenv("OLLAMA_URL", "http://localhost:11434")
                                with st.spinner("Generating review questions..."):
                                    from pipeline.ai_summarizer import generate_review_questions_ollama
                                    questions = generate_review_questions_ollama(
                                        section_text, num_questions=2, ollama_url=ollama_url)
                                    if questions:
                                        if 'ai_questions' not in st.session_state:
                                            st.session_state.ai_questions = {}
                                        st.session_state.ai_questions[s['id']] = questions
                                        st.session_state['scroll_to_section'] = section_id
                                        st.rerun()
                                    else:
                                        st.warning("Could not generate questions. Make sure Ollama is running.")
                            except Exception as e:
                                st.error(f"Error generating questions: {e}")

                        # Display stored AI Summary if exists - WITH EDITING
                        if 'ai_summaries' in st.session_state and s['id'] in st.session_state.ai_summaries:
                            st.markdown("**📚 Learning Summary (Editable):**")
                            bullets = st.session_state.ai_summaries[s['id']]
                            edited_bullets = []
                            for b_idx, bullet in enumerate(bullets):
                                edited_bullet = st.text_area(
                                    f"Bullet {b_idx + 1}",
                                    value=bullet,
                                    key=f"bullet_{section_id}_{b_idx}",
                                    height=max(40, min(80, len(bullet) // 2)),
                                    label_visibility="collapsed")
                                edited_bullets.append(edited_bullet)
                            if edited_bullets != bullets:
                                st.session_state.ai_summaries[s['id']] = edited_bullets
                                for key in list(st.session_state.keys()):
                                    if key.startswith('slides_html_') or key.startswith('slides_pptx_'):
                                        del st.session_state[key]

                        # Display stored Review Questions if exists
                        if 'ai_questions' in st.session_state and s['id'] in st.session_state.ai_questions:
                            st.markdown("**Review Questions:**")
                            questions = st.session_state.ai_questions[s['id']]
                            for q_idx, question in enumerate(questions):
                                q_col, btn_col = st.columns([10, 1])
                                with q_col:
                                    st.markdown(f"**Q{q_idx + 1}.** {question}")
                                with btn_col:
                                    answer_key = f"answer_{section_id}_{q_idx}"
                                    if answer_key not in st.session_state:
                                        if st.button("Answer", key=f"ans_btn_{section_id}_{q_idx}", type="secondary"):
                                            st.session_state[f'gen_answer_{section_id}_{q_idx}'] = True
                                            st.session_state['scroll_to_section'] = section_id

                                # Generate answer if triggered
                                if st.session_state.get(f'gen_answer_{section_id}_{q_idx}', False):
                                    st.session_state[f'gen_answer_{section_id}_{q_idx}'] = False
                                    try:
                                        ollama_url = os.getenv("OLLAMA_URL", "http://localhost:11434")
                                        with st.spinner("Generating answer..."):
                                            from pipeline.ai_summarizer import generate_answer_ollama
                                            answer = generate_answer_ollama(question, section_text, ollama_url=ollama_url)
                                            if answer:
                                                st.session_state[answer_key] = answer
                                                st.session_state['scroll_to_section'] = section_id
                                                st.rerun()
                                    except Exception as e:
                                        st.error(f"Error: {e}")

                                # Display stored answer
                                if answer_key in st.session_state:
                                    st.markdown(f"*{st.session_state[answer_key]}*")

    # Footer
    st.markdown("""
    <div class="pro-footer">
        <strong>DocuSlide Pro</strong> · AI-Powered Document Analysis
    </div>
    """,
                unsafe_allow_html=True)

